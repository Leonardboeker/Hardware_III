"""
Generate 5-slide PowerPoint proposal — same content, more space per slide.
Paper-white background, Caveat headings, Quicksand body, black & white sketch style.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

PAPER = RGBColor(0xFA, 0xFA, 0xF8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
PENCIL = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT = RGBColor(0x9A, 0x9A, 0x9E)
HIGHLIGHT_BG = RGBColor(0xE8, 0xE8, 0xE0)
RED = RGBColor(0xCC, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEADING_FONT = 'Caveat'
BODY_FONT = 'Quicksand'
TOTAL_SLIDES = 5


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PAPER


def add_text(slide, left, top, width, height, text, font_name=BODY_FONT,
             font_size=14, color=PENCIL, bold=False, italic=False,
             alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_box(slide, left, top, width, height, line_color=INK, fill_color=None, line_width=1.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.rotation = -0.3
    return shape


def add_dashed_box(slide, left, top, width, height, line_color=PENCIL):
    shape = add_box(slide, left, top, width, height, line_color=line_color)
    shape.line.dash_style = 2
    shape.line.width = Pt(1)
    return shape


def add_label(slide, left, top, text, color=INK):
    txBox = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top - 0.15),
                                      Inches(2.2), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = HEADING_FONT
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox


def footer(slide, page_num):
    add_text(slide, 0.5, 7.0, 6, 0.3,
             'Hardware III — IAAC MRAC+MAAI 2025/2026',
             font_size=8, color=LIGHT)
    add_text(slide, 11.5, 7.0, 1.5, 0.3,
             f'{page_num} / {TOTAL_SLIDES}',
             font_size=8, color=LIGHT, alignment=PP_ALIGN.RIGHT)


def slide_number(slide, num):
    add_text(slide, 11.5, 0.2, 1.5, 1.0,
             f'0{num}', font_name=HEADING_FONT, font_size=60,
             color=RGBColor(0xE8, 0xE8, 0xE0), bold=True,
             alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════
# SLIDE 1: CONCEPT & MISSION
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_number(s, 1)
footer(s, 1)

add_text(s, 0.6, 0.4, 10, 1.2,
         'Guided Comparative Assembly',
         font_name=HEADING_FONT, font_size=52, color=INK, bold=True)

add_text(s, 0.6, 1.5, 10, 0.4,
         'An interactive construction kitchen — configure, build, compare',
         font_name=BODY_FONT, font_size=16, color=LIGHT)

# Mission box
add_box(s, 0.6, 2.2, 8, 1.3, fill_color=HIGHLIGHT_BG, line_color=PENCIL)
add_label(s, 0.6, 2.2, 'Mission')
add_text(s, 0.8, 2.45, 7.6, 1.0,
         'We compare different statistics of the housing construction industry and display them through an interactive exhibit — a hands-on learning experience for children, young architects, and decision makers.',
         font_size=14, color=PENCIL)

# Method cards — larger
methods = [
    ('Masonry', 'brick by brick', 'brick or stone'),
    ('3D Printed', 'layer by layer', 'earth or concrete'),
    ('Prefab', 'module by module', 'factory grid'),
]
for i, (name, sub, mat) in enumerate(methods):
    x = 0.6 + i * 4.1
    add_box(s, x, 3.8, 3.8, 1.5)
    add_text(s, x + 0.2, 3.95, 3.4, 0.5, name,
             font_name=HEADING_FONT, font_size=28, color=INK, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 4.55, 3.4, 0.3, sub,
             font_size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 4.9, 3.4, 0.3, mat,
             font_size=12, color=LIGHT, italic=True, alignment=PP_ALIGN.CENTER)

# The Point
add_dashed_box(s, 0.6, 5.6, 8, 0.8)
add_label(s, 0.6, 5.6, 'The Point')
add_text(s, 0.8, 5.75, 7.6, 0.5,
         'Instead of reading abstract statistics, you configure your own building, watch it get constructed phase by phase, and see the real environmental and economic cost. Then build another one and compare. The data becomes physical.',
         font_size=12, color=PENCIL)

add_text(s, 0.6, 6.55, 8, 0.3,
         'Team: Leo · Elais · Rafik · Seid · Onur · Nithik',
         font_name=HEADING_FONT, font_size=13, color=LIGHT)


# ═══════════════════════════════════════════════════
# SLIDE 2: USER INTERACTION FLOW
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_number(s, 2)
footer(s, 2)

add_text(s, 0.6, 0.3, 10, 0.9,
         'User Interaction Flow',
         font_name=HEADING_FONT, font_size=44, color=INK, bold=True)
add_text(s, 0.6, 1.1, 10, 0.35,
         'The user configures everything with physical objects — no screens, no buttons',
         font_name=HEADING_FONT, font_size=13, color=RED, italic=True)

flow_steps = [
    ('1', 'Place model on\nRFID pedestal', 'Select construction\nmethod (masonry, 3DP, prefab)'),
    ('2', 'Arrange 10 pucks\non the table', 'Define building footprint\ncamera tracks distances'),
    ('3', 'Set height\n+ materials', 'Floor count marker\n+ material controller ArUco'),
    ('4', 'System validates\nfeasibility', 'Constraints enforced\n(e.g., earth 3DP max 1 story)'),
    ('5', '5 building phases\nwith animations', 'Foundation → Walls → Roof\n→ Windows → Finishing'),
    ('6', 'Save & compare\nbuilds', 'Store dataset, build another\ncompare side by side'),
]
for i, (num, title, desc) in enumerate(flow_steps):
    x = 0.4 + i * 2.1
    add_box(s, x, 1.7, 1.9, 2.8)
    add_text(s, x + 0.1, 1.8, 0.5, 0.5, num,
             font_name=HEADING_FONT, font_size=36, color=INK, bold=True)
    add_text(s, x + 0.15, 2.35, 1.6, 0.7,
             title.replace('\n', ' '),
             font_name=HEADING_FONT, font_size=14, color=INK, bold=True)
    add_text(s, x + 0.15, 3.1, 1.6, 1.0,
             desc.replace('\n', ' '),
             font_size=10, color=PENCIL)
    if i < 5:
        add_text(s, x + 1.95, 2.8, 0.2, 0.4, '→',
                 font_name=HEADING_FONT, font_size=22, color=LIGHT)

# Material constraints box
add_box(s, 0.4, 4.8, 12.5, 1.2, fill_color=HIGHLIGHT_BG, line_color=PENCIL)
add_label(s, 0.4, 4.8, 'Material Constraints')

constraints = [
    ('Masonry:', 'brick or stone'),
    ('3D Printed:', 'earth (max 1 story) or concrete'),
    ('Prefab:', 'modular grid panels'),
]
for i, (method, detail) in enumerate(constraints):
    x = 0.7 + i * 4.1
    add_text(s, x, 5.05, 1.5, 0.25, method,
             font_name=HEADING_FONT, font_size=13, color=INK, bold=True)
    add_text(s, x + 1.5, 5.05, 2.5, 0.25, detail,
             font_size=11, color=PENCIL)

add_text(s, 0.7, 5.45, 11, 0.3,
         'System shows red / error feedback when parameters are not feasible — user adjusts until valid',
         font_size=10, color=RED, italic=True)


# ═══════════════════════════════════════════════════
# SLIDE 3: SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_number(s, 3)
footer(s, 3)

add_text(s, 0.6, 0.3, 10, 0.9,
         'System Architecture',
         font_name=HEADING_FONT, font_size=44, color=INK, bold=True)
add_text(s, 0.6, 1.05, 10, 0.3,
         'Physical components → Processing → Projection feedback',
         font_size=14, color=LIGHT)

columns = [
    ('Physical Inputs', "What's on the table", [
        ('RFID pedestal + ESP', 'Place 3D model to select method'),
        ('10 ArUco pucks', 'Arrange on table for building footprint'),
        ('Height marker', 'Physically distinct ArUco for floor count'),
        ('Material controller', 'ArUco dial for foundation + material'),
        ('ESP controllers', 'Handle all hardware communication'),
    ]),
    ('Processing', 'TouchDesigner + YOLO', [
        ('Overhead camera', 'Tracks ArUco positions & distances'),
        ('YOLO plugin', 'Object recognition in TouchDesigner'),
        ('Validation engine', 'Enforces realistic constraints'),
        ('Rhino / Grasshopper', 'Geometry + parametric data'),
        ('Scaling system', 'Maps table distances to real meters'),
    ]),
    ('Output', 'Projection + Displays', [
        ('Top-down projector', 'Footprint visualization on table'),
        ('Phase animations', 'AI-generated construction sequences'),
        ('Data per phase', 'CO₂, cost, labor, time, materials'),
        ('Comparison view', 'Side-by-side multi-build dashboard'),
        ('Projection on models', '3D models fill up during phases'),
    ]),
]

for i, (label, heading, items) in enumerate(columns):
    x = 0.3 + i * 4.3
    bw = 4.1
    bh = 4.6
    add_box(s, x, 1.55, bw, bh)
    add_label(s, x, 1.55, label)

    add_text(s, x + 0.2, 1.75, bw - 0.4, 0.4, heading,
             font_name=HEADING_FONT, font_size=22, color=INK, bold=True)

    y_item = 2.25
    for title, desc in items:
        add_text(s, x + 0.2, y_item, bw - 0.4, 0.2, title,
                 font_name=BODY_FONT, font_size=12, color=INK, bold=True)
        add_text(s, x + 0.2, y_item + 0.22, bw - 0.4, 0.2, desc,
                 font_size=10, color=PENCIL)
        y_item += 0.65

# Arrows between columns
for ax_x in [4.4, 8.7]:
    add_text(s, ax_x - 0.1, 3.5, 0.3, 0.5, '→',
             font_name=HEADING_FONT, font_size=28, color=LIGHT)


# ═══════════════════════════════════════════════════
# SLIDE 4: FINITE STATE MACHINE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_number(s, 4)
footer(s, 4)

add_text(s, 0.6, 0.3, 10, 0.9,
         'Finite State Machine',
         font_name=HEADING_FONT, font_size=44, color=INK, bold=True)
add_text(s, 0.6, 1.05, 10, 0.3,
         '10 states — user configures building → system validates → guided phases → compare',
         font_size=13, color=LIGHT)

# Row 1: Setup states
setup_states = [
    ('IDLE', 'Waiting\nfor user', False),
    ('METHOD\nSELECTED', 'RFID pedestal\ntriggers', True),
    ('FOOTPRINT\nDEFINED', '10 pucks\ndefine plan', False),
    ('HEIGHT\nSET', 'Floor marker\nplaced', False),
]
setup_arrows = ['—RFID→', '—pucks placed→', '—height set→']

for i, (name, desc, filled) in enumerate(setup_states):
    x = 0.5 + i * 3.1
    fc = INK if filled else None
    tc = WHITE if filled else INK
    dc = LIGHT if filled else PENCIL
    add_box(s, x, 1.7, 2.5, 1.2, fill_color=fc if filled else None)
    add_text(s, x + 0.1, 1.78, 2.3, 0.5, name.replace('\n', ' '),
             font_name=HEADING_FONT, font_size=15, color=tc, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 2.3, 2.3, 0.5, desc.replace('\n', ' '),
             font_size=9, color=dc, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, x + 2.5, 2.1, 0.6, 0.3, setup_arrows[i],
                 font_name=HEADING_FONT, font_size=10, color=LIGHT)

# Row 2: Validation
val_states = [
    ('MATERIALS\nCHOSEN', 'Material controller\nsets options', False),
    ('VALIDATED', 'System confirms\nfeasibility', True),
    ('PHASE DISPLAY', 'Animation plays\ndata shown (×5)', False),
]
val_arrows = ['—validates→', '—begin→']

for i, (name, desc, filled) in enumerate(val_states):
    x = 0.5 + i * 3.1
    fc = INK if filled else None
    tc = WHITE if filled else INK
    dc = LIGHT if filled else PENCIL
    lw = 2.5 if i == 2 else 1.5
    add_box(s, x, 3.3, 2.5, 1.2, fill_color=fc if filled else None, line_width=lw)
    add_text(s, x + 0.1, 3.38, 2.3, 0.5, name.replace('\n', ' '),
             font_name=HEADING_FONT, font_size=15, color=tc, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 3.9, 2.3, 0.5, desc.replace('\n', ' '),
             font_size=9, color=dc, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_text(s, x + 2.5, 3.7, 0.6, 0.3, val_arrows[i],
                 font_name=HEADING_FONT, font_size=10, color=LIGHT)

# Phase display self-loop
add_text(s, 8.5, 3.3, 3.5, 0.3,
         'user accepts → next phase (loops ×5)',
         font_name=HEADING_FONT, font_size=10, color=PENCIL, italic=True)

# Error state
add_text(s, 5.2, 4.6, 1.5, 0.25, '↓ not feasible',
         font_name=HEADING_FONT, font_size=10, color=RED)
add_box(s, 4.5, 4.9, 2.5, 0.8, line_color=RED)
add_text(s, 4.6, 4.95, 2.3, 0.35, 'ERROR',
         font_name=HEADING_FONT, font_size=16, color=RED, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(s, 4.6, 5.3, 2.3, 0.3, 'Not feasible — adjust parameters',
         font_size=9, color=RED, alignment=PP_ALIGN.CENTER)
add_text(s, 7.1, 5.1, 3.5, 0.25, '—adjust→ back to MATERIALS CHOSEN',
         font_name=HEADING_FONT, font_size=10, color=RED)

# Row 3: Completion
add_text(s, 0.5, 5.95, 2, 0.25, 'all phases done →',
         font_name=HEADING_FONT, font_size=10, color=LIGHT)

add_box(s, 2.5, 5.85, 2.8, 0.8)
add_text(s, 2.6, 5.9, 2.6, 0.35, 'BUILDING COMPLETE',
         font_name=HEADING_FONT, font_size=14, color=INK, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(s, 2.6, 6.25, 2.6, 0.3, 'Final output + save dataset',
         font_size=9, color=PENCIL, alignment=PP_ALIGN.CENTER)

add_text(s, 5.35, 6.1, 0.6, 0.25, '—save→',
         font_name=HEADING_FONT, font_size=10, color=LIGHT)

add_box(s, 6.0, 5.85, 2.5, 0.8, fill_color=INK)
add_text(s, 6.1, 5.9, 2.3, 0.35, 'COMPARISON',
         font_name=HEADING_FONT, font_size=14, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(s, 6.1, 6.25, 2.3, 0.3, 'Side-by-side multi-build',
         font_size=9, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_text(s, 8.6, 6.1, 3.5, 0.25, '—build another / timeout→ IDLE',
         font_name=HEADING_FONT, font_size=10, color=LIGHT)

# Reference
add_text(s, 0.5, 6.75, 10, 0.2,
         'Ref: Stefana Parascho — Cooperative Robotic Assembly (ETH). Assembly sequence as FSM.',
         font_size=9, color=LIGHT, italic=True)


# ═══════════════════════════════════════════════════
# SLIDE 5: DATA, PRIORITIES & TIMELINE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_number(s, 5)
footer(s, 5)

add_text(s, 0.6, 0.3, 10, 0.9,
         'Data, Priorities & Timeline',
         font_name=HEADING_FONT, font_size=44, color=INK, bold=True)
add_text(s, 0.6, 1.0, 10, 0.3,
         'What we track, how we prioritize, and when we deliver',
         font_size=14, color=LIGHT)

# Parameters box — left
add_box(s, 0.4, 1.5, 6.0, 2.3)
add_label(s, 0.4, 1.5, 'Parameters per Phase')

params = [
    ('CO₂ emissions', 'Per ton of material, per phase'),
    ('Labor hours', 'Worker time + health/safety considerations'),
    ('Material usage', 'Concrete, wood, cement — contextualized (e.g., X trees)'),
    ('Total cost', 'Broken down by construction phase'),
    ('Logistics', 'Transportation, regional (Catalonia focus)'),
    ('Construction time', 'Duration per phase and total'),
]
y = 1.8
for name, desc in params:
    add_text(s, 0.6, y, 2.3, 0.2, name,
             font_name=HEADING_FONT, font_size=12, color=INK, bold=True)
    add_text(s, 2.9, y, 3.3, 0.2, desc, font_size=10, color=PENCIL)
    y += 0.32

# Priorities — right
add_box(s, 6.7, 1.5, 6.0, 2.3)
add_label(s, 6.7, 1.5, 'Development Priorities')

priorities = [
    ('P1 — Course critical', 'Projector + camera + ArUco tracking, square meter sizing, floor selection, GUI', RED),
    ('P2 — RFID + ESP', 'Pedestal integration, method selection, ESP controller communication', INK),
    ('P3 — Building phases', 'AI-generated animations, information display per construction phase', INK),
    ('P4 — Data + Comparison', 'Cost tracking, CO₂ tracking, dataset save, side-by-side comparison view', INK),
]
y = 1.8
for title, desc, color in priorities:
    add_text(s, 6.9, y, 5.5, 0.2, title,
             font_name=BODY_FONT, font_size=12, color=color, bold=True)
    add_text(s, 6.9, y + 0.22, 5.5, 0.2, desc, font_size=9, color=PENCIL)
    y += 0.52

# Building phases — left bottom
add_box(s, 0.4, 4.0, 6.0, 1.5)
add_label(s, 0.4, 4.0, 'Building Phases (~5 per method)')
phases_text = [
    '1. Foundation / Excavation',
    '2. Structure / Walls',
    '3. Roof',
    '4. Windows / Openings',
    '5. Finishing (plaster, paint, floors)',
]
y = 4.25
for pt in phases_text:
    add_text(s, 0.6, y, 5.5, 0.2, pt, font_size=11, color=PENCIL, bold=True)
    y += 0.22
add_text(s, 0.6, y + 0.05, 5.5, 0.2,
         'Exact phases per method need research — 3DP and prefab have different sequences',
         font_size=9, color=LIGHT, italic=True)

# Research action items — right bottom
add_dashed_box(s, 6.7, 4.0, 6.0, 1.5)
add_label(s, 6.7, 4.0, 'Research Action Items')
actions = [
    '☐ Research building phases for all 3 methods',
    '☐ Contact 10+ companies per method for reference data',
    '☐ Research AI animation tools for phase visualization',
    '☐ 3D print ArUco marker pucks',
    '☐ Install YOLO plugin for TouchDesigner',
    '☐ Gather Catalonia regional data (logistics, suppliers)',
]
y = 4.25
for action in actions:
    add_text(s, 6.9, y, 5.5, 0.2, action, font_size=10, color=PENCIL)
    y += 0.22

# Timeline
add_text(s, 0.6, 5.6, 4, 0.3,
         'Project Timeline',
         font_name=HEADING_FONT, font_size=22, color=INK, bold=True)

# Timeline line
line = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(1.5), Inches(6.05), Inches(10.5), Pt(2)
)
line.fill.solid()
line.fill.fore_color.rgb = INK
line.line.fill.background()

timeline_phases = [
    ('P1', 'Projector + Camera\nArUco + GUI', 'Apr 17–May 4', True),
    ('P2', 'RFID + ESP\nMethod Selection', 'May 4', False),
    ('P3', 'Building Phases\nAnimations', 'May 11', False),
    ('P4', 'Cost + CO₂\nComparison', 'May 18', False),
    ('Finals', 'Integration\nTesting + Demo', 'May 22', False),
]
for i, (label, desc, date, is_current) in enumerate(timeline_phases):
    x = 1.5 + i * 2.5
    dot = s.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.35), Inches(5.95), Inches(0.2), Inches(0.2)
    )
    if is_current:
        dot.fill.solid()
        dot.fill.fore_color.rgb = INK
    else:
        dot.fill.solid()
        dot.fill.fore_color.rgb = PAPER
    dot.line.color.rgb = INK
    dot.line.width = Pt(1.5)

    add_text(s, x, 5.72, 1.2, 0.22, label,
             font_name=HEADING_FONT, font_size=11, color=INK, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, x - 0.3, 6.15, 1.8, 0.35, desc.replace('\n', ' '),
             font_size=7.5, color=PENCIL, alignment=PP_ALIGN.CENTER)
    add_text(s, x, 6.48, 1.2, 0.18, date,
             font_name=HEADING_FONT, font_size=8, color=LIGHT,
             alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'deliverables', 'Proposal.pptx')
prs.save(output_path)
print(f"Saved {output_path}")
