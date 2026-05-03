"""Screenshot each slide of the PPTX by converting to images via LibreOffice or Playwright."""
from pptx import Presentation
from pptx.util import Inches
import subprocess, os, sys

# We'll use a workaround: export each slide as a separate HTML-like render
# Actually, let's just open the PPTX in the browser via Google Slides or use python-pptx to check dimensions

prs = Presentation(os.path.join(os.path.dirname(__file__), '..', 'deliverables', 'Proposal.pptx'))
print(f"Slides: {len(prs.slides)}")
print(f"Dimensions: {prs.slide_width / 914400:.1f}\" x {prs.slide_height / 914400:.1f}\"")

for i, slide in enumerate(prs.slides):
    shapes = list(slide.shapes)
    print(f"\nSlide {i+1}: {len(shapes)} shapes")
    for shape in shapes:
        left = shape.left / 914400
        top = shape.top / 914400
        w = shape.width / 914400
        h = shape.height / 914400
        name = shape.name[:30]
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text[:50].replace('\n', ' ')
        print(f"  {name:32s} L={left:5.1f} T={top:5.1f} W={w:4.1f} H={h:4.1f}  '{text}'")
