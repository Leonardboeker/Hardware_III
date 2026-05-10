from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Colors
BLACK    = RGBColor(0x1A, 0x1A, 0x1A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG  = RGBColor(0xF4, 0xF4, 0xF4)
ACCENT   = RGBColor(0x2A, 0x6E, 0xC8)   # blue
TERRACOTTA = RGBColor(0xC8, 0x5A, 0x1E) # masonry
GREEN    = RGBColor(0x2D, 0x8A, 0x52)   # prefab
RED_DARK = RGBColor(0x9B, 0x2D, 0x20)   # reclaimed
WARN     = RGBColor(0xD4, 0x7F, 0x00)   # warning/flag

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, BLACK)
    add_rect(slide, 0, Inches(5.8), W, Inches(1.7), ACCENT)
    add_text(slide, title,    Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.5),
             size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.5),
             size=22, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.LEFT)
    add_text(slide, "Hardware III — IAAC MRAC/MAAI 2025/2026  |  May 2026",
             Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.6),
             size=13, color=WHITE, align=PP_ALIGN.LEFT)


def section_slide(prs, label):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, ACCENT)
    add_text(slide, label, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.8),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def content_slide(prs, title, bullets, accent_color=ACCENT):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, WHITE)
    add_rect(slide, 0, 0, Inches(0.18), H, accent_color)
    add_text(slide, title, Inches(0.45), Inches(0.3), Inches(12.5), Inches(0.8),
             size=26, bold=True, color=BLACK)
    add_rect(slide, Inches(0.45), Inches(1.05), Inches(12.4), Inches(0.04), accent_color)

    y = Inches(1.25)
    for bullet in bullets:
        if bullet.startswith("##"):
            add_text(slide, bullet[2:].strip(), Inches(0.55), y,
                     Inches(12.2), Inches(0.5), size=15, bold=True, color=accent_color)
            y += Inches(0.42)
        elif bullet.startswith("!"):
            add_rect(slide, Inches(0.55), y, Inches(12.0), Inches(0.38), RGBColor(0xFF, 0xF3, 0xCD))
            add_text(slide, "⚠  " + bullet[1:].strip(), Inches(0.65), y,
                     Inches(11.8), Inches(0.38), size=13, color=WARN)
            y += Inches(0.45)
        else:
            add_text(slide, "•  " + bullet, Inches(0.65), y,
                     Inches(12.0), Inches(0.42), size=14, color=BLACK)
            y += Inches(0.42)
    return slide


def two_col_slide(prs, title, left_title, left_items, right_title, right_items,
                  accent_color=ACCENT):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, WHITE)
    add_rect(slide, 0, 0, Inches(0.18), H, accent_color)
    add_text(slide, title, Inches(0.45), Inches(0.3), Inches(12.5), Inches(0.8),
             size=26, bold=True, color=BLACK)
    add_rect(slide, Inches(0.45), Inches(1.05), Inches(12.4), Inches(0.04), accent_color)

    mid = Inches(6.85)
    add_rect(slide, mid, Inches(1.15), Inches(0.03), Inches(5.9), RGBColor(0xDD, 0xDD, 0xDD))

    add_text(slide, left_title,  Inches(0.55), Inches(1.2),  Inches(6.0), Inches(0.5),
             size=15, bold=True, color=accent_color)
    add_text(slide, right_title, Inches(7.05), Inches(1.2),  Inches(5.9), Inches(0.5),
             size=15, bold=True, color=accent_color)

    y = Inches(1.7)
    for item in left_items:
        add_text(slide, "•  " + item, Inches(0.65), y, Inches(5.9), Inches(0.45),
                 size=13, color=BLACK)
        y += Inches(0.43)

    y = Inches(1.7)
    for item in right_items:
        add_text(slide, "•  " + item, Inches(7.1), y, Inches(5.7), Inches(0.45),
                 size=13, color=BLACK)
        y += Inches(0.43)


# ── SLIDES ────────────────────────────────────────────────────────────────────

title_slide(prs,
    "Hardware III — Person 6 Role Fulfillment",
    "Data Research + Narrative Owner\nRafik El Khoury"
)

# 1 - Project Overview
content_slide(prs, "What Is Hardware III?", [
    "Interactive table installation — visitors physically configure a small building scenario",
    "Overhead projector guides assembly through 5 construction phases in real time",
    "At the end: side-by-side LCA comparison of 3 construction methods + reclaimed brick baseline",
    "Tech stack: TouchDesigner FSM + Python/OpenCV (ArUco) + ESP32/RFID",
    "Team of 7, IAAC MRAC + MAAI — finals May 22, 2026",
    "##The political argument",
    "Turn abstract construction data into embodied knowledge: you build the comparison yourself",
    "The asymmetry in data confidence between methods IS the installation's core message",
])

# 2 - Role
content_slide(prs, "Person 6 — My Role", [
    "##Responsibilities",
    "Source LCA data for all 4 construction methods",
    "Structure data as ranges with source tiers — no unsourced fixed values",
    "Write method descriptions and phase texts for visitor-facing projection",
    "Define the comparison baseline (reclaimed brick)",
    "##Deliverables",
    "4 CSV files in data/methods/ (one per method, phase × parameter × range × source)",
    "methods_db.json — corrected and populated for TouchDesigner to read",
    "Method descriptions, 5 phase texts, final comparison insight",
])

# 3 - Research Foundation
content_slide(prs, "Research Foundation — What Was Already Done", [
    "7 deep-research strands completed using the earn-the-data discipline (docs/research_2/)",
    "##Each strand produced:",
    "VALUES.md — phase × parameter × value_low/high × unit × assumption × source × tier",
    "SYNTHESIS.md — evidence convergence, methodology wobbles, geographic gaps",
    "BIBLIOGRAPHY.md — full citations, tier classification",
    "##The gap I found:",
    "Research was complete but data had never been transferred into structured files",
    "data/methods/ directory was empty — 4 CSVs and methods_db.json were all placeholders",
    "My job was not more research. It was translation: research docs → structured data",
])

# 4 - Geographic Strategy
content_slide(prs, "Geographic Strategy — Catalonia ▸ Spain ▸ EU ▸ Global", [
    "Locked in Phase 2 docs: use the most local source available, flag the level",
    "##Tier 1 (Catalan/Spanish): BEDEC (ITeC), CYPE, Hispalyt EPDs",
    "##Tier 2 (Institutional/EU): European research, validated EPDs",
    "##Tier 3 (Vendor): Allowed only with a Tier 1/2 sibling row",
    "##Per-method geographic reality:",
    "Masonry — strong Catalan data (Hispalyt, BEDEC, CYPE) ✓",
    "3D Printed — zero peer-reviewed Catalan LCA. Data: Qatar, Australia, US",
    "Prefab/CLT — no Spanish CLT mills. Data: Austrian/Swedish EPDs",
    "Reclaimed Brick — no Catalan tier-1 LCA. Data: Finland, Switzerland, UK",
    "!The geographic asymmetry between methods is not a flaw — it IS the argument",
])

# 5 - Masonry findings
content_slide(prs, "Method Finding: Masonry", [
    "Traditional fired-clay brick, hand-laid on site",
    "##CO₂ (A1–A3 cradle-to-gate)",
    "Structure: 32–75 kgCO₂e/m² wall — Hispalyt EPDs, Tier 1",
    "Whole building: 205–490 kgCO₂e/m² GFA — sanity check vs Izaola 2023, De Wolf 2017",
    "##Labour & Cost",
    "Labour: 22–34 h/m² — CYPE FEF010/FFZ010, Tier 1",
    "Cost: 950–1,350 EUR/m² — CYPE + BEDEC, Tier 1",
    "Time: 14–20 weeks",
    "##Data quality",
    "Only method with genuine Catalan tier-1 data from physical Spanish kilns",
    "BEDEC environmental column inherits Ecoinvent — not fully independent from Hispalyt EPDs",
], accent_color=TERRACOTTA)

# 6 - 3DP findings
content_slide(prs, "Method Finding: 3D Printed Concrete", [
    "Robotic concrete extrusion — structural walls only; all other phases are conventional",
    "##CO₂",
    "Wall element: 44–90 kgCO₂e/m² wall — Mohammad 2020 (Qatar), Tier 1",
    "Whole building: 58–147 kgCO₂e/m² GFA — Rossi 2024, Tier 1 (caveats: n=4)",
    "3DP earth (cob): 12–25 kgCO₂e/m² wall — Alhumayani 2020, Tier 1",
    "##7 methodology wobbles documented:",
    "Cement content (±45%), reinforcement boundary (+30%), functional unit (wall vs GFA)",
    "Printer electricity, cob biogenic convention, durability assumption, geography",
    "##Critical gap",
    "!No published LCA for any Catalan 3DP project. IAAC TOVA exists in Collserola — no LCA published.",
    "Data transfers from Qatar/Australia to Barcelona context with explicit flag",
], accent_color=ACCENT)

# 7 - Prefab findings
content_slide(prs, "Method Finding: Prefab (Modular Concrete + CLT)", [
    "Two sub-methods: volumetric modular concrete and cross-laminated timber (CLT)",
    "##CLT CO₂ (A1–A3)",
    "With biogenic carbon credit: 130–220 kgCO₂e/m² GFA",
    "Without biogenic credit: 220–350 kgCO₂e/m² GFA  ← 30% swing on accounting choice",
    "##Modular concrete CO₂",
    "280–569 kgCO₂e/m² GFA — Wei 2024 (Hong Kong typology), Tier 1",
    "##Catalonia-specific penalty most LCAs ignore:",
    "CLT Austria→Barcelona ≈ 1,800–2,200 km road = adds 25–60 kgCO₂e/m² GFA",
    "!No Spanish CLT mills exist. No Spanish modular concrete primary LCA.",
    "Data from Hong Kong, Austria, Sweden — flagged as EU proxy",
], accent_color=GREEN)

# 8 - Reclaimed Brick findings
content_slide(prs, "Method Finding: Reclaimed Brick (Baseline)", [
    "Salvaged fired-clay brick, cleaned and relaid. Role: comparison BASELINE, not a competitor",
    "##The load-bearing finding — same wall, same bricks, 3 different accounting rules:",
    "Cut-off (default):          8–25 kgCO₂e/m² GFA",
    "Avoided-burden:            18–45 kgCO₂e/m² GFA",
    "System-expansion (50/50): 35–70 kgCO₂e/m² GFA",
    "That is a 5× swing. Every published LCA makes this choice without always saying so.",
    "##Labour",
    "2–4× higher than new brick (cleaning + sorting + hand-fitting variable dimensions)",
    "##Geographic gap",
    "!No Catalan tier-1 LCA. BEDEC has no reclaimed-brick entry as of 2026.",
    "Data from Finland (Salmio & Huuhka 2026), Switzerland (Re:Crete, K.118), UK (RBC EPD)",
], accent_color=RED_DARK)

# 9 - Files delivered
two_col_slide(prs, "Files Delivered", "CSVs (data/methods/)", [
    "masonry.csv — 25 rows, Tier 1 Catalan",
    "3d-printed.csv — 33 rows, global + gap flag",
    "prefab.csv — 38 rows, CLT + modular",
    "reclaimed-brick.csv — 25 rows, 3 allocation siblings",
    "Schema: phase, parameter, value_low,",
    "  value_high, unit, assumption,",
    "  source, source_tier",
    "All UNKNOWN cells marked explicitly",
], "methods_db.json", [
    "4 correct method names (was: Cantilever/Arch/Truss)",
    "LCA summary fields per method",
    "co2_per_m2_range, labor_hours_range,",
    "  time_range, cost_per_m2_range",
    "source_label + geographic_tier",
    "Link to CSV per method",
    "5 phase texts (Foundation→Finishing)",
    "Final comparison insight (3 paragraphs)",
    "⚠ RFID tags: placeholders — Person 7 must assign",
])

# 10 - Narrative decisions
content_slide(prs, "Narrative Decisions", [
    "##Tone decision (resolved with team)",
    "Scientific, neutral, factual — not storytelling, not advocacy",
    "Each description states what the method is, where the data comes from, key number",
    "##Example — Masonry:",
    '"Load-bearing fired-clay brick, hand-laid on site. Catalan tier-1 data: Hispalyt EPDs, BEDEC (ITeC), CYPE. CO₂: 205–490 kgCO₂e/m² GFA. Labour: 22–34 h/m²."',
    "##Example — 3D Printed:",
    '"Robotic concrete extrusion. Structural walls only — foundation, roof, openings conventional. No published LCA for any Catalan project. Data from Qatar, Australia, US. CO₂: 58–147 kgCO₂e/m² GFA."',
    "##Phase texts (5 lines, one per construction phase)",
    "Functional, consistent across methods — appear during PHASE_N states in the FSM",
])

# 11 - Strength assessment
content_slide(prs, "Is This a Strong Idea? — Assessment", [
    "##STRENGTHS",
    "Earn-the-data discipline: tiered sourcing, no unsourced fixed values, ranges throughout",
    "Catalan-first geography: honest about what local data exists vs. what is a proxy",
    "The 5× reclaimed brick swing is intellectually honest and more interesting than a single number",
    "Data asymmetry between methods (masonry = real Catalan; 3DP = Qatar) IS the argument",
    "Methodology wobbles documented — the installation can show uncertainty, not false precision",
    "##CONCERNS",
    "!May 22 = 12 days. CSVs just written — not yet tested in TouchDesigner.",
    "!3 out of 4 methods have no Catalan primary LCA — honest, but will need clear UI flagging",
    "!RFID tags still placeholder — hardware dependency on Person 7 before demo",
    "Comparison view visual design not yet confirmed — data is ready; display is not",
    "##VERDICT",
    "Strong idea, methodologically serious. Risk is time, not quality of the argument.",
])

# 12 - Outstanding flags
content_slide(prs, "Outstanding Flags (Not Person 6)", [
    "##For Person 7 (Hardware/Sensors):",
    "!RFID tag IDs in methods_db.json are placeholders — read physical tags, update JSON",
    "##For the team (integration):",
    "Switch working branch to master — all Person 6 files are on treethreetree branch",
    "Test CSV reads in TouchDesigner before May 22",
    "##For the projection design:",
    "Methodology wobble overlay must surface: allocation rule for reclaimed brick,",
    "  biogenic carbon toggle for CLT, geographic source disclosure for 3DP and prefab",
    "##For Person 6 (if time allows):",
    "Strand 07 (Spanish/Catalan superpowers) was still running — check if outputs added anything",
    "Consider a SOURCES.md in data/ listing all 9 attribution corrections from data-inventory.md",
])

out = r"C:\Users\Rafik\Documents\GitHub\Hardware\deliverables\Person6_Role_Fulfillment.pptx"
import os; os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"Saved: {out}")
