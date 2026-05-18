"""Hardware III jury deck — PPTX build.

Mirrors deliverables/jury-deck/index.html. Run from repo root:
    python deliverables/jury-deck-pptx/build_pptx.py

Output: deliverables/jury-deck-pptx/HARDWARE_III_jury.pptx
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# Design tokens — match deliverables/jury-deck/assets/css/theme.css
# ============================================================

# OKLCH approximations baked into sRGB (computed externally)
PAPER        = RGBColor(0xF8, 0xF3, 0xEC)
PAPER_EDGE   = RGBColor(0xEF, 0xE9, 0xE0)
PAPER_DEEP   = RGBColor(0xE7, 0xE0, 0xD3)
RULE         = RGBColor(0xC4, 0xBB, 0xAC)
RULE_SOFT    = RGBColor(0xDB, 0xD2, 0xC2)
INK          = RGBColor(0x26, 0x21, 0x1B)
INK_STRONG   = RGBColor(0x16, 0x12, 0x0C)
INK_SOFT     = RGBColor(0x68, 0x60, 0x55)
INK_FAINT    = RGBColor(0x95, 0x8C, 0x7E)
INK_GHOST    = RGBColor(0xBE, 0xB4, 0xA4)

MASONRY      = RGBColor(0xC2, 0x6A, 0x35)
MASONRY_SOFT = RGBColor(0xF1, 0xDE, 0xD0)
PRINTED      = RGBColor(0x3D, 0x82, 0xC2)
PRINTED_SOFT = RGBColor(0xCF, 0xDE, 0xEE)
PREFAB       = RGBColor(0x3F, 0xA8, 0x6E)
PREFAB_SOFT  = RGBColor(0xCF, 0xE8, 0xD9)

LIVE         = RGBColor(0x3F, 0xB1, 0x70)
PENDING      = RGBColor(0xD8, 0x9F, 0x3E)
INVALID      = RGBColor(0xC4, 0x3D, 0x2E)

# Fonts — picks fall back gracefully if not installed
F_DISPLAY = "Funnel Display"
F_SANS    = "Funnel Sans"
F_MONO    = "JetBrains Mono"
F_DISPLAY_FALLBACK = "Calibri"
F_SANS_FALLBACK    = "Calibri"
F_MONO_FALLBACK    = "Consolas"

# Slide canvas — 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Grid: kicker row, body, foot row
MARGIN_X = Inches(0.55)
KICKER_Y = Inches(0.35)
KICKER_H = Inches(0.32)
FOOT_Y   = Inches(7.0)
FOOT_H   = Inches(0.32)

# Body sits between kicker and foot with breathing room
BODY_Y = Inches(0.95)
BODY_H = Inches(5.85)


# ============================================================
# Helpers
# ============================================================

def _i(v):
    """Coerce EMU-ish value to int. python-pptx connector serialization
    chokes on floats, dropping ".0" suffix into OOXML attributes which
    PowerPoint refuses to open. All x/y/w/h that come from arithmetic
    (Inches division, EMU averaging) must pass through this."""
    return int(round(float(v)))


def add_text_box(slide, x, y, w, h, text, *, font=F_SANS, fallback=F_SANS_FALLBACK,
                 size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.15,
                 letter_spacing=None, uppercase=False, space_after=0):
    """Add a single-paragraph text box."""
    box = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    if uppercase and isinstance(text, str):
        text = text.upper()

    paragraphs = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if letter_spacing:
            _set_letter_spacing(run, letter_spacing)
    return box


def _set_letter_spacing(run, hundredths_of_point):
    """Apply tracking via the OOXML 'spc' attribute. Units are hundredths of a point.
    100 = 1pt of tracking. Use ~80-160 for uppercase labels."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_of_point)))


def add_rect(slide, x, y, w, h, *, fill=PAPER, line=None, line_width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), _i(w), _i(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    # No text inside this rect; ensure no autofill text
    shape.text_frame.text = ""
    return shape


def add_hline(slide, x, y, w, *, color=RULE_SOFT, weight=0.75):
    line = slide.shapes.add_connector(1, _i(x), _i(y), _i(x + w), _i(y))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_kicker(slide, slide_no, section, suffix=None):
    text = f"{slide_no:02d} / 19   ·   {section.upper()}"
    if suffix:
        text += f"   ·   {suffix}"
    box = add_text_box(
        slide, MARGIN_X, KICKER_Y, SLIDE_W - 2 * MARGIN_X, KICKER_H,
        text, font=F_MONO, fallback=F_MONO_FALLBACK,
        size=10, color=INK_FAINT, letter_spacing=120,
    )
    # Highlight section label in stronger ink — rewrite as 3 runs
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    def _add(text, color=INK_FAINT, bold=False):
        r = p.add_run()
        r.text = text
        r.font.name = F_MONO
        r.font.size = Pt(10)
        r.font.color.rgb = color
        r.font.bold = bold
        _set_letter_spacing(r, 120)
    _add(f"{slide_no:02d} / 19", color=INK_FAINT)
    _add("   ·   ", color=INK_FAINT)
    _add(section.upper(), color=INK_STRONG, bold=True)
    if suffix:
        _add("   ·   ", color=INK_FAINT)
        _add(suffix.upper(), color=INK_FAINT)
    return box


def add_foot(slide, slide_no, section):
    # Top rule
    add_hline(slide, MARGIN_X, FOOT_Y, SLIDE_W - 2 * MARGIN_X, color=RULE_SOFT)
    # Left brand mark
    add_text_box(
        slide, MARGIN_X, FOOT_Y + Inches(0.07), Inches(6), FOOT_H,
        "HARDWARE III · GUIDED COMPARATIVE ASSEMBLY",
        font=F_MONO, fallback=F_MONO_FALLBACK,
        size=8, color=INK_FAINT, letter_spacing=140,
    )
    # Right slide number
    add_text_box(
        slide, SLIDE_W - MARGIN_X - Inches(3.5), FOOT_Y + Inches(0.07),
        Inches(3.5), FOOT_H,
        f"{slide_no:02d} / 19 · {section.upper()}",
        font=F_MONO, fallback=F_MONO_FALLBACK,
        size=8, color=INK_SOFT, letter_spacing=140, align=PP_ALIGN.RIGHT,
    )


def add_paper_background(slide):
    # Solid paper-colored background rectangle behind everything
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Push to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_paper_background(slide)
    return slide


def add_eyebrow(slide, x, y, w, text):
    return add_text_box(
        slide, x, y, w, Inches(0.3), text.upper(),
        font=F_MONO, fallback=F_MONO_FALLBACK,
        size=9, color=INK_SOFT, bold=True, letter_spacing=180,
    )


def add_chip(slide, x, y, w, h, *, n, name, sub, highlight=False, name_size=14):
    fill = MASONRY_SOFT if highlight else PAPER
    line = MASONRY if highlight else RULE
    add_rect(slide, x, y, w, h, fill=fill, line=line, line_width=0.75)
    # number
    add_text_box(
        slide, x + Inches(0.12), y + Inches(0.08), w - Inches(0.2), Inches(0.22),
        n, font=F_MONO, fallback=F_MONO_FALLBACK,
        size=8, color=INK_FAINT, letter_spacing=120,
    )
    # name — let it fill the middle, can wrap to two lines
    add_text_box(
        slide, x + Inches(0.12), y + Inches(0.32), w - Inches(0.16), h - Inches(0.6),
        name, font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=name_size, bold=True, color=INK_STRONG, line_spacing=1.05,
    )
    # sub
    if sub:
        add_text_box(
            slide, x + Inches(0.12), y + h - Inches(0.28), w - Inches(0.2), Inches(0.24),
            sub, font=F_MONO, fallback=F_MONO_FALLBACK,
            size=7.5, color=INK_SOFT,
        )


def add_tag(slide, x, y, w, h, label):
    add_rect(slide, x, y, w, h, fill=PAPER, line=RULE, line_width=0.75)
    add_text_box(
        slide, x + Inches(0.1), y, w - Inches(0.2), h,
        label.upper(), font=F_MONO, fallback=F_MONO_FALLBACK,
        size=8, color=INK, bold=True, letter_spacing=140,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_box_node(slide, x, y, w, h, *, title, subtitle, detail, line_width=1.0, accent=False):
    fill = MASONRY_SOFT if accent else PAPER
    line = MASONRY if accent else INK_STRONG
    add_rect(slide, x, y, w, h, fill=fill, line=line, line_width=line_width)
    add_text_box(
        slide, x + Inches(0.14), y + Inches(0.1), w - Inches(0.28), Inches(0.32),
        title.upper(), font=F_MONO, fallback=F_MONO_FALLBACK,
        size=10, bold=True, color=INK_STRONG, letter_spacing=140,
    )
    add_text_box(
        slide, x + Inches(0.14), y + Inches(0.42), w - Inches(0.28), Inches(0.28),
        subtitle, font=F_SANS, fallback=F_SANS_FALLBACK,
        size=9, color=INK, bold=True,
    )
    if detail:
        add_text_box(
            slide, x + Inches(0.14), y + Inches(0.7), w - Inches(0.28), h - Inches(0.78),
            detail, font=F_MONO, fallback=F_MONO_FALLBACK,
            size=8, color=INK_SOFT, letter_spacing=80,
        )


def add_arrow(slide, x1, y1, x2, y2, *, color=INK_STRONG, weight=1.5):
    line = slide.shapes.add_connector(1, _i(x1), _i(y1), _i(x2), _i(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # Add arrow head
    lnEl = line.line._get_or_add_ln()
    tail = lnEl.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(lnEl, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def add_bullet(slide, x, y, w, h, items, *, size=11, color=INK, font=F_SANS, fallback=F_SANS_FALLBACK, line_spacing=1.4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        marker = p.add_run()
        marker.text = "+ "
        marker.font.name = F_MONO
        marker.font.size = Pt(size)
        marker.font.color.rgb = INK_FAINT
        marker.font.bold = True
        body = p.add_run()
        body.text = item
        body.font.name = font
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def add_spec(slide, x, y, w, items, *, size=9):
    """items = [(label, value), ...]"""
    row_h = Inches(0.26)
    label_w = w * 0.42
    value_w = w - label_w
    for i, (label, value) in enumerate(items):
        ry = y + row_h * i
        add_text_box(
            slide, x, ry, label_w, row_h,
            label.upper(), font=F_MONO, fallback=F_MONO_FALLBACK,
            size=size - 1, color=INK_FAINT, letter_spacing=140,
        )
        add_text_box(
            slide, x + label_w, ry, value_w, row_h,
            value, font=F_MONO, fallback=F_MONO_FALLBACK,
            size=size, color=INK, bold=True,
        )


def add_table_rect(slide, x, y, col_w, rows, *, header):
    """Simple table with header + body. col_w is list of widths in Inches."""
    total_w = sum(col_w)
    # Header rule above
    add_hline(slide, x, y, total_w, color=RULE, weight=1.0)
    row_h = Inches(0.36)
    # Header
    cx = x
    for i, label in enumerate(header):
        add_text_box(
            slide, cx + Inches(0.08), y + Inches(0.08), col_w[i] - Inches(0.16), row_h,
            label.upper(), font=F_MONO, fallback=F_MONO_FALLBACK,
            size=8, color=INK_FAINT, bold=True, letter_spacing=140,
        )
        cx += col_w[i]
    add_hline(slide, x, y + row_h, total_w, color=RULE, weight=0.75)
    # Body rows
    body_y = y + row_h + Inches(0.04)
    for row in rows:
        cx = x
        rh = Inches(0.5)  # auto-grows visually if text wraps
        for i, val in enumerate(row):
            add_text_box(
                slide, cx + Inches(0.08), body_y + Inches(0.06), col_w[i] - Inches(0.16), rh,
                val, font=F_MONO if i == 0 else F_SANS,
                fallback=F_MONO_FALLBACK if i == 0 else F_SANS_FALLBACK,
                size=8.5, color=INK if i == 0 else INK_SOFT,
                bold=(i == 0),
            )
            cx += col_w[i]
        add_hline(slide, x, body_y + rh, total_w, color=RULE_SOFT, weight=0.5)
        body_y += rh


# ============================================================
# Slide builders
# ============================================================

def slide_00_title(prs):
    s = new_slide(prs)

    add_text_box(
        s, MARGIN_X, Inches(0.55), SLIDE_W - 2 * MARGIN_X, Inches(0.32),
        "IAAC   ·   HARDWARE III   ·   HUMAN IN THE LOOP · INTERACTIVE SYSTEMS   ·   MRAC + MAAI 2025-2026",
        font=F_MONO, fallback=F_MONO_FALLBACK,
        size=9, color=INK_FAINT, letter_spacing=160,
    )
    add_text_box(
        s, MARGIN_X, Inches(1.7), SLIDE_W - 2 * MARGIN_X, Inches(4.5),
        "guided\ncomparative\nassembly",
        font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=88, bold=True, color=INK_STRONG, line_spacing=0.95,
    )
    # Meta row
    meta_y = Inches(6.05)
    col_w = Inches(2.9)
    metas = [
        ("Team", "Leo, Elais, Rafik, Seid, Onur, Nithik"),
        ("Schedule", "Apr 10 to May 22, 2026"),
        ("Final critique", "May 22, 2026"),
        ("Runtime", "TouchDesigner + Python + ESP32"),
    ]
    for i, (label, value) in enumerate(metas):
        x = MARGIN_X + col_w * i
        add_text_box(s, x, meta_y, col_w, Inches(0.22), label.upper(),
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK_FAINT, bold=True, letter_spacing=180)
        add_text_box(s, x, meta_y + Inches(0.28), col_w, Inches(0.4), value,
                     font=F_SANS, fallback=F_SANS_FALLBACK,
                     size=10, color=INK)

    add_foot(s, 0, "Title")
    return s


def slide_01_context_problem(prs):
    s = new_slide(prs)
    add_kicker(s, 1, "Context", "the problem")

    add_eyebrow(s, MARGIN_X, BODY_Y, Inches(4), "The problem")
    add_text_box(
        s, MARGIN_X, BODY_Y + Inches(0.4), Inches(6.2), Inches(3.5),
        "construction decisions\ncarry hidden\nconsequences",
        font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=44, bold=True, color=INK_STRONG, line_spacing=0.98,
    )

    # Right column
    rx = Inches(7.2)
    rw = Inches(5.6)
    add_text_box(
        s, rx, BODY_Y + Inches(0.4), rw, Inches(1.5),
        "Material selection, construction method, footprint and height are decisions made early in a project, often before environmental or economic data is consulted.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=12, color=INK, line_spacing=1.5,
    )
    add_text_box(
        s, rx, BODY_Y + Inches(1.95), rw, Inches(1.5),
        "For non experts this data is inaccessible. For experts it is contested. Methodological assumptions vary and single figure comparisons mislead.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=12, color=INK, line_spacing=1.5,
    )
    add_hline(s, rx, BODY_Y + Inches(3.6), rw, color=RULE, weight=0.75)
    add_text_box(
        s, rx, BODY_Y + Inches(3.85), rw, Inches(1.4),
        "What if you could see those consequences as you configure the building itself?",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=16, color=INK_STRONG, bold=True, line_spacing=1.35,
    )

    add_foot(s, 1, "Context")
    return s


def slide_02_context_proposition(prs, video_path):
    s = new_slide(prs)
    add_kicker(s, 2, "Context", "the proposition")

    # Left: video / poster
    fig_x = MARGIN_X
    fig_y = BODY_Y
    fig_w = Inches(7.5)
    fig_h = Inches(4.7)
    # background panel
    add_rect(s, fig_x, fig_y, fig_w, fig_h, fill=PAPER_EDGE, line=RULE, line_width=0.75)
    if video_path and os.path.exists(video_path):
        poster = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                              "assets", "demo_poster.jpg")
        try:
            kwargs = dict(mime_type="video/mp4")
            if os.path.exists(poster):
                kwargs["poster_frame_image"] = poster
            s.shapes.add_movie(
                video_path,
                _i(fig_x + Inches(0.04)), _i(fig_y + Inches(0.04)),
                _i(fig_w - Inches(0.08)), _i(fig_h - Inches(0.08)),
                **kwargs,
            )
        except Exception as e:
            add_text_box(
                s, fig_x, fig_y, fig_w, fig_h,
                f"[ video embed failed: {e} ]",
                font=F_MONO, fallback=F_MONO_FALLBACK,
                size=10, color=INK_FAINT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
    else:
        add_text_box(
            s, fig_x, fig_y, fig_w, fig_h,
            "[ demo video file missing ]",
            font=F_MONO, fallback=F_MONO_FALLBACK,
            size=10, color=INK_FAINT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    # Live pill bottom-left
    pill_x = fig_x + Inches(0.2)
    pill_y = fig_y + fig_h - Inches(0.5)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_x, pill_y, Inches(1.9), Inches(0.32))
    pill.fill.solid()
    pill.fill.fore_color.rgb = PAPER
    pill.line.color.rgb = RULE
    pill.line.width = Pt(0.75)
    pill.shadow.inherit = False
    add_text_box(
        s, pill_x + Inches(0.18), pill_y, Inches(1.7), Inches(0.32),
        "● LIVE · DETECTION DEMO",
        font=F_MONO, fallback=F_MONO_FALLBACK,
        size=7.5, color=INK, bold=True, letter_spacing=140,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Right column
    rx = Inches(8.4)
    rw = Inches(4.4)
    add_text_box(
        s, rx, BODY_Y, rw, Inches(1.9),
        "a table\nthat reads\nyour building",
        font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=34, bold=True, color=INK_STRONG, line_spacing=0.98,
    )
    add_text_box(
        s, rx, BODY_Y + Inches(2.05), rw, Inches(2.0),
        "An overhead short-throw projector returns visual feedback on a 9-panel layout. Configuration happens on the table itself. The system reads, computes, and shows the consequences live.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=11, color=INK, line_spacing=1.5,
    )
    add_spec(s, rx, BODY_Y + Inches(4.1), rw, [
        ("Projection", "1280 x 720, short throw"),
        ("Surface", "table-mounted"),
        ("Loop", "closed via overhead webcam"),
    ])

    add_foot(s, 2, "Context")
    return s


def slide_03_experience(prs):
    s = new_slide(prs)
    add_kicker(s, 3, "The experience")

    # Left
    add_text_box(
        s, MARGIN_X, BODY_Y, Inches(7.2), Inches(2.7),
        "a table that reads\nyour configuration,\nand projects its\nconsequences",
        font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=42, bold=True, color=INK_STRONG, line_spacing=0.98,
    )
    add_bullet(s, MARGIN_X, BODY_Y + Inches(3.0), Inches(7.2), Inches(2.8), [
        "Place a physical method model on the RFID pedestal to choose the construction method.",
        "Set the working plane with four ArUco corner markers on the table.",
        "Dwell a red puck three seconds at each point to draw the footprint polygon, projection responds live.",
        "Set height with a slider, set material with an ArUco marker.",
        "Walk through five construction phases with projected data per phase. End on a comparison view.",
    ], size=10)

    # Right
    rx = Inches(8.4)
    rw = Inches(4.4)
    add_eyebrow(s, rx, BODY_Y, rw, "Output per phase")
    output = [
        ("CO₂ eq", "range", "kg, sourced"),
        ("Embodied energy", "range", "MJ, sourced"),
        ("Labour", "range", "hours, sourced"),
        ("Construction time", "range", "days, sourced"),
        ("Cost", "range", "EUR 2026, sourced"),
    ]
    ry = BODY_Y + Inches(0.42)
    for label, value, unit in output:
        add_text_box(s, rx, ry, Inches(1.8), Inches(0.3), label.upper(),
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8.5, color=INK_FAINT, letter_spacing=120)
        add_text_box(s, rx + Inches(1.8), ry, Inches(1.0), Inches(0.3), value,
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=9, color=INK, bold=True)
        add_text_box(s, rx + Inches(2.8), ry, Inches(1.6), Inches(0.3), unit,
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8.5, color=INK_SOFT)
        ry += Inches(0.32)
        add_hline(s, rx, ry, rw, color=RULE_SOFT, weight=0.5)
        ry += Inches(0.04)
    add_text_box(
        s, rx, ry + Inches(0.2), rw, Inches(1.6),
        "Shown as sourced ranges with provenance, not single figures. Methodological wobble (assumption choice, geometry efficiency, biogenic carbon accounting) is a first class layer, not a footnote.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=9.5, color=INK_SOFT, line_spacing=1.5,
    )

    add_foot(s, 3, "Experience")
    return s


def slide_04_interaction(prs):
    s = new_slide(prs)
    add_kicker(s, 4, "Interaction", "8 steps · every transition confirmed by the camera")

    # 8 chips
    chip_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.42)) / 8  # 7 gaps
    chip_h = Inches(1.6)
    chip_y = BODY_Y + Inches(0.6)
    chips = [
        ("01", "Attract", "idle animation"),
        ("02", "Method", "RFID select"),
        ("03", "Footprint", "10 pucks"),
        ("04", "Height", "slider"),
        ("05", "Materials", "marker"),
        ("06", "Validate", "summary"),
        ("07", "5 phases", "data per phase"),
        ("08", "Compare", "final view"),
    ]
    # Setup / Configuration / Assembly / Output labels
    band_y = BODY_Y + Inches(0.2)
    bands = ["SETUP", "CONFIGURATION", "ASSEMBLY", "OUTPUT"]
    band_w = (SLIDE_W - 2 * MARGIN_X) / 4
    for i, label in enumerate(bands):
        add_text_box(s, MARGIN_X + band_w * i, band_y, band_w, Inches(0.3),
                     label, font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK_FAINT, letter_spacing=180, align=PP_ALIGN.CENTER)
    for i, (n, name, sub) in enumerate(chips):
        x = MARGIN_X + (chip_w + Inches(0.06)) * i
        add_chip(s, x, chip_y, chip_w, chip_h, n=n, name=name, sub=sub, highlight=(i == 6))

    # Two parallel tracks: phase mode (masonry, 3DP) + lifecycle mode (prefab)
    tracks_y = chip_y + chip_h + Inches(0.35)
    half_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) / 2
    sub_w = (half_w - Inches(0.32)) / 5
    sub_h = Inches(1.2)

    add_eyebrow(s, MARGIN_X, tracks_y, half_w, "Phase mode · masonry · 3D printed")
    for i, (n, name) in enumerate([
        ("P1", "Foundation"),
        ("P2", "Walls"),
        ("P3", "Roof"),
        ("P4", "Openings"),
        ("P5", "Finishing"),
    ]):
        x = MARGIN_X + (sub_w + Inches(0.08)) * i
        add_chip(s, x, tracks_y + Inches(0.32), sub_w, sub_h, n=n, name=name, sub="", name_size=12)

    lx = MARGIN_X + half_w + Inches(0.4)
    add_eyebrow(s, lx, tracks_y, half_w, "Lifecycle mode · prefab · EN 15978")
    for i, (n, name) in enumerate([
        ("A1-A3", "Product"),
        ("A4", "Transport"),
        ("A5", "Assembly"),
        ("B", "Use"),
        ("C", "End of life"),
    ]):
        x = lx + (sub_w + Inches(0.08)) * i
        add_chip(s, x, tracks_y + Inches(0.32), sub_w, sub_h, n=n, name=name, sub="", name_size=12)

    add_foot(s, 4, "Interaction")
    return s


def slide_05_logic(prs):
    s = new_slide(prs)
    add_kicker(s, 5, "Logic architecture", "three layers, separated on purpose")

    # Content FSM row
    add_eyebrow(s, MARGIN_X, BODY_Y, Inches(10), "Content FSM · visitor facing · 8 states")
    chip_y = BODY_Y + Inches(0.4)
    chip_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.42)) / 8
    chip_h = Inches(1.45)
    states = [
        ("01", "Idle", "attract"),
        ("02", "Method", "RFID select"),
        ("03", "Footprint", "10 pucks"),
        ("04", "Height", "floor slider"),
        ("05", "Materials", "material marker"),
        ("06", "Validated", "summary"),
        ("07", "Phase N", "1 to 5"),
        ("08", "Compare", "paired metrics"),
    ]
    for i, (n, name, sub) in enumerate(states):
        x = MARGIN_X + (chip_w + Inches(0.06)) * i
        add_chip(s, x, chip_y, chip_w, chip_h, n=n, name=name, sub=sub, highlight=(i == 6))

    # Wrapper + Visual sub-layers
    sub_y = chip_y + chip_h + Inches(0.45)
    half_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) / 2

    add_eyebrow(s, MARGIN_X, sub_y, half_w, "Wrapper · setup & recovery · interrupts content FSM")
    wrapper_tags = ["Calibration check", "Error", "Reset", "Manual override"]
    tx = MARGIN_X
    ty = sub_y + Inches(0.4)
    for tag in wrapper_tags:
        tw = Inches(1.55)
        add_tag(s, tx, ty, tw, Inches(0.36), tag)
        tx += tw + Inches(0.08)

    vx = MARGIN_X + half_w + Inches(0.4)
    add_eyebrow(s, vx, sub_y, half_w, "Visual feedback codes · projection output, not FSM states")
    visual_tags = ["Disconnected", "Pending", "Invalid", "Valid", "Idle anim", "Summary", "Comparison"]
    tx = vx
    ty = sub_y + Inches(0.4)
    for tag in visual_tags:
        tw = Inches(0.95)
        add_tag(s, tx, ty, tw, Inches(0.36), tag)
        tx += tw + Inches(0.06)

    add_foot(s, 5, "Logic")
    return s


def slide_06_architecture(prs):
    s = new_slide(prs)
    add_kicker(s, 6, "Architecture", "input · runtime · output")

    # 3 lanes
    lane_x = MARGIN_X
    lane_w = SLIDE_W - 2 * MARGIN_X
    lane_h = Inches(1.4)
    gap    = Inches(0.18)
    lane_y = BODY_Y + Inches(0.3)

    for i, (label, fill) in enumerate([("INPUT", PAPER_EDGE), ("RUNTIME", PAPER_DEEP), ("OUTPUT", PAPER_EDGE)]):
        y = lane_y + (lane_h + gap) * i
        add_rect(s, lane_x, y, lane_w, lane_h, fill=fill, line=None)
        add_text_box(s, lane_x + Inches(0.18), y + Inches(0.1), Inches(2), Inches(0.3),
                     label, font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK_FAINT, letter_spacing=180)

    # Input boxes
    box_h = Inches(1.0)
    box_w = Inches(2.5)
    iy = lane_y + Inches(0.3)
    add_box_node(s, MARGIN_X + Inches(1.8),  iy, box_w, box_h,
                 title="ESP32", subtitle="+ MFRC522 RFID",
                 detail="USB serial · method id")
    add_box_node(s, MARGIN_X + Inches(5.4), iy, box_w, box_h,
                 title="USB Webcam", subtitle="overhead, 1080p",
                 detail="locked focus + exposure")
    add_box_node(s, MARGIN_X + Inches(9.0),  iy, Inches(2.9), box_h,
                 title="Python + OpenCV", subtitle="ArUco detection · homography",
                 detail="OSC out, port 7000")

    # Runtime
    ry = lane_y + lane_h + gap + Inches(0.3)
    add_box_node(s, SLIDE_W / 2 - Inches(1.85), ry, Inches(3.7), box_h,
                 title="TouchDesigner", subtitle="primary runtime",
                 detail="FSM · compute · render · output", line_width=1.75)

    # Output
    oy = lane_y + (lane_h + gap) * 2 + Inches(0.3)
    add_box_node(s, SLIDE_W / 2 - Inches(1.85), oy, Inches(3.7), box_h,
                 title="Short throw projector", subtitle="1280 x 720, table surface",
                 detail="homography aligned")

    # Arrows
    add_arrow(s, MARGIN_X + Inches(3.0), iy + box_h,
              SLIDE_W / 2 - Inches(1.1), ry)
    add_arrow(s, MARGIN_X + Inches(10.4), iy + box_h,
              SLIDE_W / 2 + Inches(1.1), ry)
    add_arrow(s, SLIDE_W / 2, ry + box_h,
              SLIDE_W / 2, oy)
    # Webcam → OpenCV horizontal arrow
    add_arrow(s, MARGIN_X + Inches(6.65), iy + box_h / 2,
              MARGIN_X + Inches(9.0), iy + box_h / 2)

    add_foot(s, 6, "Architecture")
    return s


def slide_07_vision(prs):
    s = new_slide(prs)
    add_kicker(s, 7, "Computer vision", "USB webcam, 1080p, overhead")

    # 5 boxes
    pipe = [
        ("Capture", "USB webcam\n1080p overhead", False),
        ("ArUco detect", "Python · OpenCV 4.8+\n4 x 4 dictionary per frame", False),
        ("Homography", "cam px → projector px\n4 calibration YAML, one time", False),
        ("Grace period", "10 frames\ntimestamp + stale detection", False),
        ("OSC out", "TouchDesigner\nUDP port 7000", True),
    ]
    bw = (SLIDE_W - 2 * MARGIN_X - Inches(0.6)) / 5
    by = BODY_Y + Inches(0.2)
    bh = Inches(1.7)
    for i, (title, detail, terminal) in enumerate(pipe):
        x = MARGIN_X + (bw + Inches(0.15)) * i
        accent = terminal
        add_rect(s, x, by, bw, bh,
                 fill=MASONRY_SOFT if accent else PAPER,
                 line=MASONRY if accent else RULE, line_width=0.75)
        add_text_box(s, x + Inches(0.18), by + Inches(0.18), bw - Inches(0.36), Inches(0.6),
                     title, font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
                     size=18, bold=True, color=INK_STRONG, line_spacing=1.0)
        add_text_box(s, x + Inches(0.18), by + Inches(0.85), bw - Inches(0.36), bh - Inches(1.0),
                     detail, font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK_SOFT, line_spacing=1.4, letter_spacing=80)

    # OSC channel notes
    cy = by + bh + Inches(0.55)
    col_w = (SLIDE_W - 2 * MARGIN_X) / 3 - Inches(0.2)
    channels = [
        ("/puck/N i f f",
         "Per detected puck: frame number, projector_x, projector_y. Sent every frame the puck is visible."),
        ("/puck/lost i",
         "Sent when a puck leaves the camera frame. Carries the puck ID."),
        ("/vision/heartbeat i",
         "Every frame, current frame number. Puck data goes stale after 10 frames without a heartbeat. Triggers Error via wrapper FSM."),
    ]
    for i, (label, body) in enumerate(channels):
        x = MARGIN_X + (col_w + Inches(0.3)) * i
        add_text_box(s, x, cy, col_w, Inches(0.3), label,
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=10, color=INK_STRONG, bold=True, letter_spacing=120)
        add_text_box(s, x, cy + Inches(0.4), col_w, Inches(1.6), body,
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8.5, color=INK_SOFT, line_spacing=1.55)

    add_foot(s, 7, "Vision")
    return s


def slide_08_gestures(prs):
    s = new_slide(prs)
    add_kicker(s, 8, "Hand gesture sketch", "vision2 · MediaPipe Tasks API")

    # Left
    add_text_box(
        s, MARGIN_X, BODY_Y, Inches(6.5), Inches(1.6),
        "drawing the\nfootprint by hand",
        font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=32, bold=True, color=INK_STRONG, line_spacing=0.98,
    )
    add_text_box(
        s, MARGIN_X, BODY_Y + Inches(1.8), Inches(6.5), Inches(1.7),
        "ArUco IDs 0 to 3 fix the working plane via homography. A red puck inside the plane places footprint corner points. MediaPipe Tasks reads the hand for secondary actions, anywhere in the camera frame.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=11, color=INK, line_spacing=1.5,
    )
    add_spec(s, MARGIN_X, BODY_Y + Inches(3.7), Inches(6.5), [
        ("Working plane", "ArUco IDs 0 to 3, 900 × 600 px"),
        ("Method cards", "ArUco IDs 20 · 21 · 22"),
        ("Puck", "HSV red blob, dwell timer"),
        ("Hand detector", "MediaPipe Tasks · hand_landmarker.task"),
    ], size=9)

    # Right: gestures
    rx = Inches(7.6)
    rw = Inches(5.2)
    add_eyebrow(s, rx, BODY_Y, rw, "Puck · primary placement")
    # primary puck card
    puck_y = BODY_Y + Inches(0.4)
    add_rect(s, rx, puck_y, rw, Inches(1.0), fill=PAPER, line=INVALID, line_width=0.75)
    add_text_box(s, rx + Inches(0.18), puck_y + Inches(0.1), rw - Inches(0.36), Inches(0.35),
                 "place point", font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
                 size=16, bold=True, color=INK_STRONG)
    add_text_box(s, rx + Inches(0.18), puck_y + Inches(0.45), rw - Inches(0.36), Inches(0.25),
                 "RED PUCK INSIDE THE WORKING PLANE",
                 font=F_MONO, fallback=F_MONO_FALLBACK,
                 size=8, color=INK_SOFT, letter_spacing=120)
    add_text_box(s, rx + Inches(0.18), puck_y + Inches(0.7), rw - Inches(0.36), Inches(0.3),
                 "hold still 3 s · re-arm by moving away and back",
                 font=F_MONO, fallback=F_MONO_FALLBACK,
                 size=8, color=INK_FAINT)

    add_eyebrow(s, rx, puck_y + Inches(1.2), rw, "Hand gestures · secondary actions")
    gx = rx
    gy = puck_y + Inches(1.6)
    cw = (rw - Inches(0.15)) / 2
    ch = Inches(1.1)
    gests = [
        ("undo", "FLAT FIST · PALM DOWN", "hold 2 s · removes last point"),
        ("add window", "PEACE · V · INDEX + MIDDLE UP", "hold 2 s · nearest wall"),
        ("extrude", "INDEX FINGER UP", "hold 3 s · walls into 3D"),
        ("reset", "UPRIGHT FIST · VERTICAL", "hold 5 s · clears the sketch"),
    ]
    for i, (name, glyph, hold) in enumerate(gests):
        col = i % 2
        row = i // 2
        cx = gx + (cw + Inches(0.15)) * col
        cy_ = gy + (ch + Inches(0.15)) * row
        add_rect(s, cx, cy_, cw, ch, fill=PAPER, line=RULE, line_width=0.6)
        add_text_box(s, cx + Inches(0.14), cy_ + Inches(0.1), cw - Inches(0.28), Inches(0.32),
                     name, font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
                     size=14, bold=True, color=INK_STRONG)
        add_text_box(s, cx + Inches(0.14), cy_ + Inches(0.45), cw - Inches(0.28), Inches(0.3),
                     glyph, font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=7.5, color=INK_SOFT, letter_spacing=80)
        add_text_box(s, cx + Inches(0.14), cy_ + Inches(0.75), cw - Inches(0.28), Inches(0.3),
                     hold, font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=7.5, color=INK_FAINT)

    add_foot(s, 8, "Sketch")
    return s


def slide_09_touchdesigner(prs):
    s = new_slide(prs)
    add_kicker(s, 9, "TouchDesigner runtime", "seven canonical nodes · the vertical slice network")

    box_w = Inches(2.4)
    box_h = Inches(1.05)

    # 2 inputs left
    add_box_node(s, MARGIN_X + Inches(0.2), BODY_Y + Inches(0.5),
                 box_w, box_h,
                 title="vision_in", subtitle="OSC In CHOP",
                 detail="puck positions + heartbeat from CV")
    add_box_node(s, MARGIN_X + Inches(0.2), BODY_Y + Inches(3.7),
                 box_w, box_h,
                 title="rfid_in", subtitle="Constant CHOP / Serial DAT",
                 detail="method selector")

    # Compute_state center
    cs_x = MARGIN_X + Inches(3.6)
    cs_y = BODY_Y + Inches(2.05)
    cs_w = Inches(2.8)
    cs_h = Inches(1.4)
    add_rect(s, cs_x, cs_y, cs_w, cs_h, fill=PAPER, line=INK_STRONG, line_width=1.5)
    add_text_box(s, cs_x + Inches(0.18), cs_y + Inches(0.1), cs_w - Inches(0.36), Inches(0.32),
                 "compute_state".upper(), font=F_MONO, fallback=F_MONO_FALLBACK,
                 size=10, bold=True, color=INK_STRONG, letter_spacing=140)
    add_text_box(s, cs_x + Inches(0.18), cs_y + Inches(0.4), cs_w - Inches(0.36), Inches(0.3),
                 "Script CHOP · every frame", font=F_SANS, fallback=F_SANS_FALLBACK,
                 size=10, color=INK, bold=True)
    add_text_box(s, cs_x + Inches(0.18), cs_y + Inches(0.7), cs_w - Inches(0.36), Inches(0.6),
                 "puck_count · area_px2\nmethod_id · hb_alive",
                 font=F_MONO, fallback=F_MONO_FALLBACK,
                 size=8.5, color=INK_SOFT, line_spacing=1.4)

    # 2 outputs right column
    rx = MARGIN_X + Inches(7.4)
    add_box_node(s, rx, BODY_Y + Inches(0.5), box_w, box_h,
                 title="render_footprint", subtitle="Script TOP · 1280 x 720",
                 detail="polygon · puck circles · color band")
    add_box_node(s, rx, BODY_Y + Inches(3.7), box_w, box_h,
                 title="stats_text", subtitle="Text TOP",
                 detail="puck count · area · heartbeat")

    # projector_out final
    px = MARGIN_X + Inches(10.4)
    py = BODY_Y + Inches(2.05)
    add_box_node(s, px, py, Inches(2.3), cs_h,
                 title="projector_out", subtitle="Window COMP",
                 detail="full screen output")

    # arrows
    add_arrow(s, MARGIN_X + Inches(2.6), BODY_Y + Inches(1.03),  cs_x, cs_y + Inches(0.5))
    add_arrow(s, MARGIN_X + Inches(2.6), BODY_Y + Inches(4.23),  cs_x, cs_y + cs_h - Inches(0.5))
    add_arrow(s, cs_x + cs_w, cs_y + Inches(0.5), rx, BODY_Y + Inches(1.03))
    add_arrow(s, cs_x + cs_w, cs_y + cs_h - Inches(0.5), rx, BODY_Y + Inches(4.23))
    add_arrow(s, rx + box_w, BODY_Y + Inches(1.03), px, py + Inches(0.5))
    add_arrow(s, rx + box_w, BODY_Y + Inches(4.23), px, py + cs_h - Inches(0.5))

    add_foot(s, 9, "TouchDesigner")
    return s


def slide_10_physical(prs):
    s = new_slide(prs)
    add_kicker(s, 10, "Physical layer", "objects the visitor interacts with")

    cols = 3
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.5)) / cols
    cx0 = MARGIN_X
    cy0 = BODY_Y
    icon_h = Inches(2.0)

    def col(i, count, title, bullets, icon_factory):
        x = cx0 + (col_w + Inches(0.25)) * i
        add_rect(s, x, cy0, col_w, icon_h, fill=PAPER_EDGE, line=RULE, line_width=0.5)
        icon_factory(x, cy0, col_w, icon_h)
        add_text_box(s, x, cy0 + icon_h + Inches(0.08), col_w, Inches(0.25),
                     count, font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK_FAINT, letter_spacing=160)
        add_text_box(s, x, cy0 + icon_h + Inches(0.32), col_w, Inches(0.55),
                     title, font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
                     size=20, bold=True, color=INK_STRONG, line_spacing=1.0)
        add_bullet(s, x, cy0 + icon_h + Inches(0.95), col_w, Inches(2.6), bullets, size=9, line_spacing=1.4)

    def icon_rfid(x, y, w, h):
        gx = x + w / 2 - Inches(0.9)
        gy = y + h / 2 - Inches(0.55)
        for k in range(2):
            xx = gx + Inches(1.0) * k
            add_rect(s, xx, gy, Inches(0.85), Inches(1.05), fill=PAPER, line=INK, line_width=1.0)
            add_text_box(s, xx, gy + Inches(0.3), Inches(0.85), Inches(0.5),
                         "RFID", font=F_MONO, fallback=F_MONO_FALLBACK,
                         size=8, color=INK, bold=True, letter_spacing=120, align=PP_ALIGN.CENTER)

    def icon_plane_puck(x, y, w, h):
        # 4 corners
        sz = Inches(0.5)
        margin = Inches(0.4)
        positions = [(x + margin, y + margin),
                     (x + w - margin - sz, y + margin),
                     (x + margin, y + h - margin - sz),
                     (x + w - margin - sz, y + h - margin - sz)]
        for px, py in positions:
            add_rect(s, px, py, sz, sz, fill=PAPER, line=INK, line_width=1.0)
        # dashed plane outline simulated as a rectangle
        plane = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            x + margin + sz + Inches(0.05), y + margin + sz + Inches(0.05),
            w - 2 * (margin + sz) - Inches(0.1), h - 2 * (margin + sz) - Inches(0.1))
        plane.fill.background()
        plane.line.color.rgb = RULE
        plane.line.width = Pt(0.5)
        plane.line.dash_style = 7  # DASH
        plane.shadow.inherit = False
        # red puck centered
        cx = x + w / 2
        cy = y + h / 2
        puck = s.shapes.add_shape(MSO_SHAPE.OVAL,
            cx - Inches(0.25), cy - Inches(0.25), Inches(0.5), Inches(0.5))
        puck.fill.solid()
        puck.fill.fore_color.rgb = INVALID
        puck.line.color.rgb = INK
        puck.line.width = Pt(1.0)
        puck.shadow.inherit = False

    def icon_inputs(x, y, w, h):
        sx = x + w / 2 - Inches(1.3)
        sy = y + h / 2 - Inches(0.25)
        # slider
        add_rect(s, sx, sy, Inches(1.2), Inches(0.3), fill=PAPER, line=INK, line_width=1.0)
        add_rect(s, sx + Inches(0.45), sy - Inches(0.1), Inches(0.25), Inches(0.5),
                 fill=PAPER, line=INK, line_width=1.0)
        add_text_box(s, sx, sy + Inches(0.5), Inches(1.2), Inches(0.3),
                     "Height", font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK, letter_spacing=120, align=PP_ALIGN.CENTER)
        # material
        mx = sx + Inches(1.5)
        add_rect(s, mx, sy - Inches(0.1), Inches(1.0), Inches(0.5), fill=PAPER, line=INK, line_width=1.0)
        add_rect(s, mx + Inches(0.18), sy + Inches(0.0), Inches(0.65), Inches(0.3), fill=RULE_SOFT, line=None)
        add_text_box(s, mx, sy + Inches(0.5), Inches(1.0), Inches(0.3),
                     "Material", font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK, letter_spacing=120, align=PP_ALIGN.CENTER)

    col(0, "x3", "method models",
        ["3D printed physical representations of masonry, 3D printed and prefab construction methods.",
         "Each contains a MIFARE classic 1k RFID tag.",
         "Placed on the RFID pedestal to trigger Method state."],
        icon_rfid)
    col(1, "4 + 1", "plane + footprint puck",
        ["ArUco IDs 0 to 3 tape the 4 table corners and fix the working plane via homography.",
         "One red puck (HSV tracked) dwells 3 seconds inside the plane to place a footprint point.",
         "Repeat ten times. Polygon area via shoelace formula."],
        icon_plane_puck)
    col(2, "x2", "configuration inputs",
        ["Height: DollaTek 10K linear slide potentiometer wired to ESP32 GPIO34, ADC into the existing serial pipe. Replaces ArUco dial ID 10.",
         "Material: ArUco ID 11, same detection pipeline as footprint pucks, placement gated by FSM."],
        icon_inputs)

    add_foot(s, 10, "Physical layer")
    return s


def slide_11_slider(prs):
    s = new_slide(prs)
    add_kicker(s, 11, "Height slider · Phase 02.1", "DollaTek 10K linear-slide potentiometer")

    # Left
    add_text_box(
        s, MARGIN_X, BODY_Y, Inches(6.5), Inches(1.9),
        "a physical\nslider for floors,\nnot a marker",
        font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=36, bold=True, color=INK_STRONG, line_spacing=0.98,
    )
    add_text_box(
        s, MARGIN_X, BODY_Y + Inches(2.1), Inches(6.5), Inches(1.6),
        "The original ArUco dial for height (ID 10) was retired. A DollaTek 10K linear-slide potentiometer drives an ADC on the existing ESP32 RFID board, applies median + EMA smoothing and floor-boundary hysteresis, and emits the result over the same Serial pipe.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=11, color=INK, line_spacing=1.5,
    )
    add_text_box(
        s, MARGIN_X, BODY_Y + Inches(3.85), Inches(6.5), Inches(1.5),
        "HEIGHT now reads as a continuous physical action, not a discrete card placement. Per method floor caps (3D printed earth = 1, CLT = 8, modular concrete = 12) are honored from methods_db.json.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=11, color=INK_SOFT, line_spacing=1.5,
    )

    # Right: pipeline diagram
    rx = Inches(7.6)
    rw = Inches(5.2)
    bx = rx
    by = BODY_Y + Inches(0.2)
    bw = Inches(2.3)
    bh = Inches(0.7)
    add_box_node(s, bx, by, bw, bh,
                 title="Slider 10K", subtitle="DollaTek linear", detail="")
    add_box_node(s, bx + Inches(2.6), by, bw, bh,
                 title="ESP32 ADC", subtitle="GPIO34 · median + EMA", detail="")
    add_arrow(s, bx + bw, by + bh / 2, bx + Inches(2.6), by + bh / 2)
    add_box_node(s, bx + Inches(2.6), by + Inches(1.0), bw, bh,
                 title="USB Serial", subtitle="115200 · FLOOR:N · SLIDER:0.xxx", detail="")
    add_arrow(s, bx + Inches(2.6) + bw / 2, by + bh,
              bx + Inches(2.6) + bw / 2, by + Inches(1.0))
    add_box_node(s, bx + Inches(2.6), by + Inches(2.0), bw, bh,
                 title="TouchDesigner", subtitle="serial_rfid_v1 → floor", detail="")
    add_arrow(s, bx + Inches(2.6) + bw / 2, by + Inches(1.0) + bh,
              bx + Inches(2.6) + bw / 2, by + Inches(2.0))
    add_box_node(s, bx + Inches(0.0), by + Inches(2.0), Inches(1.6), bh,
                 title="Height", subtitle="FSM", detail="", line_width=1.75)
    add_arrow(s, bx + Inches(2.6), by + Inches(2.0) + bh / 2,
              bx + Inches(1.6), by + Inches(2.0) + bh / 2)

    # Spec
    add_spec(s, rx, by + Inches(3.2), rw, [
        ("Floor cap", "per method, methods_db.json"),
        ("Smoothing", "median + EMA + hysteresis"),
        ("Latency", "≤ 100 ms slider to projection"),
    ], size=9)

    add_foot(s, 11, "Slider")
    return s


def slide_12_data(prs):
    s = new_slide(prs)
    add_kicker(s, 12, "Data layer", "three methods · sourced ranges · known wobble")

    headers = ["", "Method", "Description", "Primary sources", "Known wobble", "Metric per phase"]
    rows = [
        ["", "Masonry", "Traditional fired clay block, Catalonia reference.",
         "CYPE · BEDEC (ITeC) · EPD international",
         "Regional variation in clay firing energy intensity.",
         "CO₂ · Energy · Labor · Time · Cost"],
        ["", "3D printed", "Concrete extrusion. ICON, Apis Cor, COBOD.",
         "Journal papers · vendor EPDs (tier 2 / 3)",
         "Cement content range · geometry efficiency factor.",
         "CO₂ · Energy · Labor · Time · Cost"],
        ["", "Prefab", "Modular concrete (MiC) and mass timber CLT.",
         "EU JRC · EPFL SXL · institutional reports",
         "Biogenic carbon accounting · reuse allocation rules.",
         "CO₂ · Energy · Labor · Time · Cost"],
    ]

    # Build a real PPTX table for crispness
    rows_n = 1 + len(rows)
    cols_n = len(headers)
    table_w = SLIDE_W - 2 * MARGIN_X
    table_h = Inches(3.3)
    shape = s.shapes.add_table(rows_n, cols_n, MARGIN_X, BODY_Y, table_w, table_h)
    tbl = shape.table

    col_widths = [0.45, 1.4, 2.7, 2.7, 2.7, 2.3]  # sum ~12.25
    factor = (table_w / Inches(1)) / sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w * factor)
    tbl.rows[0].height = Inches(0.45)
    for r in range(1, rows_n):
        tbl.rows[r].height = Inches(0.95)

    # Header row
    for i, label in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = label.upper()
        run.font.name = F_MONO
        run.font.size = Pt(8)
        run.font.color.rgb = INK_FAINT
        run.font.bold = True
        _set_letter_spacing(run, 160)
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.04)

    # Body rows
    method_swatches = [MASONRY, PRINTED, PREFAB]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
            if c == 0:
                # color swatch
                cell.fill.solid()
                cell.fill.fore_color.rgb = method_swatches[r]
                cell.text = ""
            else:
                cell.text = ""
                para = cell.text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.LEFT
                run = para.add_run()
                run.text = val
                run.font.name = F_MONO if c == 1 else F_SANS
                run.font.size = Pt(9)
                run.font.color.rgb = INK if c == 1 else INK_SOFT
                run.font.bold = (c == 1)

    # Footnote
    add_text_box(
        s, MARGIN_X, BODY_Y + Inches(3.5), SLIDE_W - 2 * MARGIN_X, Inches(2.0),
        "Every figure carries a confidence level. Numbers are stored as ranges with assumption notes so the methodology-wobble layer can draw from them. DOIs for all Tier 1 citations are listed in data/SOURCES.md (verified 2026-05-18).",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=10, color=INK_SOFT, line_spacing=1.5,
    )

    add_foot(s, 12, "Data")
    return s


def slide_13_metrics(prs):
    s = new_slide(prs)
    add_kicker(s, 13, "Metrics engine", "data → UI state · one source of truth")

    # 4 boxes in a row + a panels row below
    by = BODY_Y + Inches(0.4)
    bw = Inches(2.8)
    bh = Inches(1.2)
    gap = Inches(0.3)
    boxes = [
        ("methods_db.json", "4 method CSVs", "45+ source citations", False),
        ("metrics/pipeline", "normalize · range", "shape factor · units", False),
        ("metrics_engine", "scenario in, metric out", "phase or lifecycle mode", True),
        ("ui_state", "panel-ready strings", "guidance · highlights", False),
    ]
    total_w = bw * 4 + gap * 3
    bx0 = (SLIDE_W - total_w) / 2
    for i, (title, sub, det, accent) in enumerate(boxes):
        x = bx0 + (bw + gap) * i
        add_box_node(s, x, by, bw, bh, title=title, subtitle=sub, detail=det,
                     line_width=1.75 if accent else 1.0, accent=accent)
        if i < 3:
            add_arrow(s, x + bw, by + bh / 2, x + bw + gap, by + bh / 2)

    # Panels row
    panels_y = by + bh + Inches(0.7)
    panels_w = bw * 3 + gap * 2
    panels_x = bx0 + bw + gap
    add_rect(s, panels_x, panels_y, panels_w, Inches(0.85),
             fill=PAPER, line=RULE, line_width=0.75)
    add_text_box(s, panels_x + Inches(0.18), panels_y + Inches(0.12),
                 panels_w - Inches(0.36), Inches(0.32),
                 "9 PROJECTION PANELS",
                 font=F_MONO, fallback=F_MONO_FALLBACK,
                 size=10, bold=True, color=INK_STRONG, letter_spacing=140)
    add_text_box(s, panels_x + Inches(0.18), panels_y + Inches(0.42),
                 panels_w - Inches(0.36), Inches(0.4),
                 "text_top_phase_navigation · text_left_info · text_right_comparison · text_right_cost_chart · text_right_phase_preview · ...",
                 font=F_MONO, fallback=F_MONO_FALLBACK,
                 size=7.5, color=INK_SOFT)
    # Down arrow from ui_state box
    add_arrow(s, bx0 + (bw + gap) * 3 + bw / 2, by + bh,
              panels_x + panels_w * 0.85, panels_y)

    # 3 mode notes below
    nx0 = MARGIN_X
    nw = (SLIDE_W - 2 * MARGIN_X - Inches(0.6)) / 3
    ny = panels_y + Inches(1.4)
    notes = [
        ("Phase mode", "Foundation, walls, roof, openings, finishing. Used by masonry and 3D printed."),
        ("Lifecycle mode", "A1-A3 · A4 · A5 · B · C. Used by prefab (CLT and modular concrete sub-methods)."),
        ("Guidance", "The engine emits the next user action explicitly. The status bar reads it."),
    ]
    for i, (label, body) in enumerate(notes):
        nx = nx0 + (nw + Inches(0.3)) * i
        add_eyebrow(s, nx, ny, nw, label)
        add_text_box(s, nx, ny + Inches(0.4), nw, Inches(1.3),
                     body, font=F_SANS, fallback=F_SANS_FALLBACK,
                     size=10, color=INK_SOFT, line_spacing=1.5)

    add_foot(s, 13, "Metrics")
    return s


def _panel_mock(s, x, y, w, h, *, labeled=False, annotated=False, role_colors=None):
    """Draw the 9-panel layout mock at (x,y,w,h)."""
    add_rect(s, x, y, w, h, fill=PAPER_DEEP, line=RULE, line_width=0.75)
    gap = Inches(0.07)
    pad = Inches(0.08)
    inner_x = x + pad
    inner_y = y + pad
    inner_w = w - 2 * pad
    inner_h = h - 2 * pad
    # 3 columns: 1 / 1.8 / 1
    col_ratio = (1.0, 1.8, 1.0)
    cs = sum(col_ratio)
    col_widths = [inner_w * r / cs - gap * 2 / 3 for r in col_ratio]
    # 4 rows + status bar (0.22 ratio)
    row_ratio = (1.0, 1.0, 1.0, 1.0, 0.22)
    rs = sum(row_ratio)
    row_total_h = inner_h - gap * 4
    row_heights = [row_total_h * r / rs for r in row_ratio]
    # Compute positions
    cxs = [inner_x]
    for i in range(2):
        cxs.append(cxs[-1] + col_widths[i] + gap)
    rys = [inner_y]
    for i in range(4):
        rys.append(rys[-1] + row_heights[i] + gap)

    def cell(area, c, r, rowspan=1, role="output", label="", role_text=""):
        cx = cxs[c]
        cy = rys[r]
        cw = col_widths[c]
        ch = sum(row_heights[r:r + rowspan]) + gap * (rowspan - 1)
        fill = PAPER_EDGE
        if role_colors:
            fill = role_colors.get(role, PAPER_EDGE)
        add_rect(s, cx, cy, cw, ch, fill=fill, line=RULE_SOFT, line_width=0.5)
        if labeled or annotated:
            add_text_box(s, cx + Inches(0.06), cy + Inches(0.04), cw - Inches(0.12), Inches(0.2),
                         label, font=F_MONO, fallback=F_MONO_FALLBACK,
                         size=7, bold=True, color=INK_STRONG, letter_spacing=140)
            if role_text:
                add_text_box(s, cx + Inches(0.06), cy + ch - Inches(0.22), cw - Inches(0.12), Inches(0.18),
                             role_text, font=F_MONO, fallback=F_MONO_FALLBACK,
                             size=6, color=INK_FAINT, letter_spacing=100)

    # Areas
    cell("info",   0, 0, rowspan=3, role="input",  label="INFO",            role_text="method · floors · area")
    cell("phase",  1, 0, rowspan=1, role="input",  label="PHASE NAVIGATION", role_text="← zone geometry")
    cell("part",   2, 0, rowspan=2, role="output", label="PART IMPACT",     role_text="← metric engine")
    cell("plan",   1, 1, rowspan=2, role="input",  label="PLAN INTERACTION", role_text="10 pucks → footprint")
    cell("total",  2, 2, rowspan=1, role="output", label="TOTAL IMPACT",    role_text="whole-building totals")
    cell("assy",   0, 3, rowspan=1, role="output", label="ASSEMBLY SEQUENCE", role_text="phase summaries")
    cell("method", 1, 3, rowspan=1, role="input",  label="METHOD SELECTION", role_text="RFID active")
    cell("state",  2, 3, rowspan=1, role="state",  label="STATE PREVIEW",   role_text="workflow / UI state")

    # Status bar
    sx = inner_x
    sw = inner_w
    sy = rys[4]
    sh = row_heights[4]
    add_rect(s, sx, sy, sw, sh, fill=PAPER_EDGE, line=RULE_SOFT, line_width=0.5)
    if labeled or annotated:
        add_text_box(s, sx + Inches(0.08), sy, sw, sh,
                     "STATUS BAR · VISION LIVE · NEXT ACTION" if annotated else "status bar",
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=7, color=INK_FAINT, letter_spacing=140, anchor=MSO_ANCHOR.MIDDLE)


def slide_14_layout_empty(prs):
    s = new_slide(prs)
    add_kicker(s, 14, "Layout", "9-panel projection · 1280 x 720")
    mw = Inches(9.5)
    mh = Inches(4.4)
    mx = (SLIDE_W - mw) / 2
    my = BODY_Y + Inches(0.3)
    _panel_mock(s, mx, my, mw, mh)
    add_text_box(
        s, MARGIN_X, my + mh + Inches(0.4), SLIDE_W - 2 * MARGIN_X, Inches(1.2),
        "The projection is composed as one Script TOP that lays out nine named panels. Each panel has a stable role and a stable position. The layout is the visual contract.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=11, color=INK_SOFT, line_spacing=1.5, align=PP_ALIGN.CENTER,
    )
    add_foot(s, 14, "Layout")
    return s


def slide_15_layout_roles(prs):
    s = new_slide(prs)
    add_kicker(s, 15, "Layout", "three panel roles")
    mw = Inches(9.5)
    mh = Inches(4.0)
    mx = (SLIDE_W - mw) / 2
    my = BODY_Y + Inches(0.2)
    # Role-colored fills
    input_c  = RGBColor(0xC2, 0xB6, 0xA0)
    output_c = PAPER_EDGE
    state_c  = PAPER
    _panel_mock(s, mx, my, mw, mh, labeled=True,
                role_colors={"input": input_c, "output": output_c, "state": state_c})

    # Legend
    ly = my + mh + Inches(0.4)
    legend_x = MARGIN_X + Inches(0.5)
    items = [
        ("Input driven", "accepts user inputs or external data", input_c),
        ("Output driven", "presents computed data or results", output_c),
        ("State reflector", "reflects current FSM state", state_c),
    ]
    for i, (label, sub, color) in enumerate(items):
        x = legend_x + Inches(4.0) * i
        sw_x = x
        sw = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sw_x, ly, Inches(0.3), Inches(0.3))
        sw.fill.solid()
        sw.fill.fore_color.rgb = color
        sw.line.color.rgb = RULE
        sw.line.width = Pt(0.5)
        sw.shadow.inherit = False
        add_text_box(s, x + Inches(0.45), ly - Inches(0.04), Inches(3.4), Inches(0.3),
                     label, font=F_SANS, fallback=F_SANS_FALLBACK,
                     size=11, bold=True, color=INK_STRONG)
        add_text_box(s, x + Inches(0.45), ly + Inches(0.2), Inches(3.4), Inches(0.3),
                     sub, font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8.5, color=INK_SOFT)
    add_foot(s, 15, "Layout")
    return s


def slide_16_layout_annotated(prs):
    s = new_slide(prs)
    add_kicker(s, 16, "Layout", "nine named panels · data sources annotated")
    mw = Inches(8.2)
    mh = Inches(5.0)
    mx = MARGIN_X
    my = BODY_Y
    input_c  = RGBColor(0xC2, 0xB6, 0xA0)
    output_c = PAPER_EDGE
    state_c  = PAPER
    _panel_mock(s, mx, my, mw, mh, annotated=True,
                role_colors={"input": input_c, "output": output_c, "state": state_c})

    # Reading order panel right
    rx = mx + mw + Inches(0.4)
    rw = SLIDE_W - rx - MARGIN_X
    add_eyebrow(s, rx, my, rw, "Reading order")
    add_text_box(
        s, rx, my + Inches(0.4), rw, Inches(2.0),
        "Top-left clusters the visitor's selection. The center is the plan. The right column reads impacts.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=11, color=INK, line_spacing=1.55,
    )
    add_text_box(
        s, rx, my + Inches(2.5), rw, Inches(2.5),
        "The bottom row carries assembly history, method selection, and current state. The status bar runs along the floor of the projection.",
        font=F_SANS, fallback=F_SANS_FALLBACK,
        size=11, color=INK_SOFT, line_spacing=1.55,
    )

    add_foot(s, 16, "Layout")
    return s


def slide_17_fabrication(prs):
    s = new_slide(prs)
    add_kicker(s, 17, "Fabrication", "components and references")

    headers = ["Component", "Quantity", "Method", "Notes"]
    rows = [
        ["Method models", "3", "FDM print + Rhino model",
         "Masonry, 3D printed, prefab. Building + 3D print phase 1-4 FBX. Rhino sources on Juan branch."],
        ["Plane corner markers", "4", "A3 print, 4 x 4 dictionary",
         "ArUco IDs 0 to 3. cad/aruco-markers/PRINT_A3.pdf."],
        ["Footprint puck", "1", "Red object, HSV tracked",
         "Dwell-to-place point. Detector in vision/vision2/puck_detector.py."],
        ["Method ID markers", "3", "A3 print, 4 x 4 dictionary",
         "ArUco IDs 20 · 21 · 22 for masonry / 3D printed / prefab."],
        ["Height + material markers", "2", "A3 print, 4 x 4 dictionary",
         "ArUco ID 10 height fallback, ID 11 material. Slider supersedes ID 10."],
        ["RFID pedestal", "1", "FDM enclosure",
         "cad/enclosure/render_enclosure.py + render. ESP32 + RC522 wiring documented."],
        ["ESP32 + RC522 firmware", "1", "Arduino sketch",
         "methods_db.json drives tag mapping. Lab walkthrough committed."],
        ["Camera + projector rig", "1", "Overhead frame, MX Brio",
         "Focus, exposure, white balance locked."],
    ]

    rows_n = 1 + len(rows)
    cols_n = len(headers)
    table_w = SLIDE_W - 2 * MARGIN_X
    table_h = Inches(5.5)
    shape = s.shapes.add_table(rows_n, cols_n, MARGIN_X, BODY_Y, table_w, table_h)
    tbl = shape.table
    col_widths = [2.6, 1.0, 2.6, 6.05]
    factor = (table_w / Inches(1)) / sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w * factor)
    tbl.rows[0].height = Inches(0.4)
    for r in range(1, rows_n):
        tbl.rows[r].height = Inches(0.62)

    for i, label in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = label.upper()
        run.font.name = F_MONO
        run.font.size = Pt(8)
        run.font.color.rgb = INK_FAINT
        run.font.bold = True
        _set_letter_spacing(run, 160)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = val
            run.font.name = F_MONO if c == 0 else F_SANS
            run.font.size = Pt(8.5)
            run.font.color.rgb = INK if c == 0 else INK_SOFT
            run.font.bold = (c == 0)

    add_foot(s, 17, "Fabrication")
    return s


BIB_3DP = [
    ("Alhumayani, H., Gomaa, M., Soebarto, V., Jabi, W. (2020).",
     "Environmental assessment of large-scale 3D printing in construction: A comparative study between cob and concrete.",
     "Journal of Cleaner Production, 270, 122463.",
     "DOI 10.1016/j.jclepro.2020.122463"),
    ("Allouzi, R., Al-Azhari, W., Allouzi, R. (2020).",
     "Conventional Construction and 3D Printing: A Comparison Study on Material Cost in Jordan.",
     "Journal of Engineering, 2020, 1424682.",
     "DOI 10.1155/2020/1424682"),
    ("Mohammed, M., Rahman, R., Mohamed, S. F., Ahmad, M. (2020).",
     "3D Concrete Printing Sustainability: A Comparative Life Cycle Assessment of Four Construction Method Scenarios.",
     "Buildings, 10(12), 245.",
     "DOI 10.3390/buildings10120245"),
    ("Motalebi, A., Khondoker, M. A. H., Kabir, G. (2024).",
     "A systematic review of life cycle assessments of 3D concrete printing.",
     "Sustainable Operations and Computers, 5, 41–50.",
     "sciencedirect.com/science/article/pii/S2666412723000132"),
    ("Rossi, B., et al. (2024).",
     "Comparison of Embodied Carbon of 3D-printed vs. Conventionally Built Houses.",
     "Findings, Feb 2024.",
     "DOI 10.32866/001c.89707"),
]

BIB_MAS_PRE = [
    ("Izaola, B., Akizu-Gardoki, O., Oregi, X. (2023).",
     "Setting baselines of the embodied, operational and whole life carbon emissions of the average Spanish residential building.",
     "Sustainable Production and Consumption, 40, 252–264.",
     "sciencedirect.com/science/article/pii/S2352550923001574"),
    ("Ferreira, A., Pinheiro, M. D., de Brito, J., Mateus, R. (2023).",
     "Embodied vs. Operational Energy and Carbon in Retail Building Shells: A Case Study in Portugal.",
     "Energies, 16(1), 378.",
     "DOI 10.3390/en16010378"),
    ("Andersen, J. H., Rasmussen, N. L., Ryberg, M. W. (2022).",
     "Comparative life cycle assessment of cross laminated timber building and concrete building with special focus on biogenic carbon.",
     "Energy and Buildings, 254, 111604.",
     "sciencedirect.com/science/article/pii/S0378778821008884"),
    ("Hemmati, M., Messadi, T., Gu, H., Seddelmeyer, J. (2024).",
     "Comparison of Embodied Carbon Footprint of a Mass Timber Building Structure with a Steel Equivalent.",
     "Buildings, 14(5), 1276.",
     "DOI 10.3390/buildings14051276"),
    ("Wei, J., Ge, B., Zhong, Y., et al. (2024).",
     "Comparative analysis of embodied carbon in modular and conventional construction methods in Hong Kong.",
     "Scientific Reports, 14, 22833.",
     "DOI 10.1038/s41598-024-73906-7"),
]

BIB_DATA = [
    ("ITeC Banco BEDEC, 2025/2026 release.",
     "Catalonia construction cost and material-quantity database.",
     "itec.cat/nouBedec.c/bedec.aspx", ""),
    ("CYPE Arquímedes · Generador de Precios.",
     "Spanish cost-reference lines (FEF010, FFZ010) and labour rates.",
     "generadordeprecios.info", ""),
    ("Hispalyt GlobalEPD 008-001 / 008-016 / 008-017.",
     "Spanish clay roof-tile and fired-clay brick environmental product declarations.",
     "hispalyt.es / globalepd.org", ""),
    ("Stora Enso · KLH · Binderholz EPDs.",
     "Cross-laminated timber panel environmental product declarations used for biogenic carbon and end-of-life triangulation.",
     "via EPD Norway · IBU · manufacturers' registries", ""),
    ("EN 15978 + EN 15804 + A2.",
     "European standards for life-cycle assessment of buildings and EPD methodology with biogenic carbon module D.",
     "cen.eu", ""),
]


def _bib_entry(slide, x, y, w, item):
    """One bibliography entry. Author + title + venue + DOI on four short lines."""
    author, title, venue, doi = item
    cy = y
    add_text_box(slide, x, cy, w, Inches(0.24), author,
                 font=F_SANS, fallback=F_SANS_FALLBACK,
                 size=8.5, bold=True, color=INK_STRONG, line_spacing=1.25)
    cy += Inches(0.22)
    add_text_box(slide, x, cy, w, Inches(0.38), title,
                 font=F_SANS, fallback=F_SANS_FALLBACK,
                 size=8.5, italic=True, color=INK, line_spacing=1.3)
    cy += Inches(0.34)
    add_text_box(slide, x, cy, w, Inches(0.22), venue,
                 font=F_SANS, fallback=F_SANS_FALLBACK,
                 size=8, color=INK_SOFT, line_spacing=1.25)
    cy += Inches(0.20)
    if doi:
        add_text_box(slide, x, cy, w, Inches(0.22), doi,
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=7.5, color=INK_SOFT, letter_spacing=80)


def slide_18_bibliography(prs):
    s = new_slide(prs)
    add_kicker(s, 18, "Bibliography", "tier 1 academic sources · databases · standards")

    # 3 columns
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.6)) / 3
    col_x = [
        MARGIN_X,
        MARGIN_X + col_w + Inches(0.3),
        MARGIN_X + (col_w + Inches(0.3)) * 2,
    ]
    headers = [
        "3D-printed concrete + earth",
        "Masonry + Prefab",
        "Databases · EPDs · Standards",
    ]
    items_by_col = [BIB_3DP, BIB_MAS_PRE, BIB_DATA]

    for c in range(3):
        add_eyebrow(s, col_x[c], BODY_Y, col_w, headers[c])
        add_hline(s, col_x[c], BODY_Y + Inches(0.32), col_w, color=RULE)
        entry_y = BODY_Y + Inches(0.42)
        for item in items_by_col[c]:
            _bib_entry(s, col_x[c], entry_y, col_w, item)
            entry_y += Inches(1.08)

    add_foot(s, 18, "Bibliography")
    return s


def slide_19_end(prs):
    s = new_slide(prs)
    add_text_box(
        s, MARGIN_X, KICKER_Y, SLIDE_W - 2 * MARGIN_X, KICKER_H,
        "19 / 19   ·   THANK YOU",
        font=F_MONO, fallback=F_MONO_FALLBACK,
        size=10, color=INK_FAINT, bold=True, letter_spacing=160,
    )
    add_text_box(
        s, MARGIN_X, Inches(1.8), SLIDE_W - 2 * MARGIN_X, Inches(3.6),
        "questions\nwelcome",
        font=F_DISPLAY, fallback=F_DISPLAY_FALLBACK,
        size=110, bold=True, color=INK_STRONG, line_spacing=0.95,
    )
    add_hline(s, MARGIN_X, Inches(5.9), SLIDE_W - 2 * MARGIN_X, color=RULE_SOFT)
    metas = [
        ("Repository", "github.com/Leonardboeker/Hardware_III"),
        ("Demo", "vertical-slice.toe"),
        ("Team", "Leo, Elais, Rafik, Seid, Onur, Nithik"),
    ]
    meta_y = Inches(6.15)
    col_w = Inches(4.0)
    for i, (label, value) in enumerate(metas):
        x = MARGIN_X + col_w * i
        add_text_box(s, x, meta_y, col_w, Inches(0.22), label.upper(),
                     font=F_MONO, fallback=F_MONO_FALLBACK,
                     size=8, color=INK_FAINT, bold=True, letter_spacing=180)
        add_text_box(s, x, meta_y + Inches(0.28), col_w, Inches(0.4), value,
                     font=F_SANS, fallback=F_SANS_FALLBACK,
                     size=10, color=INK)
    add_foot(s, 19, "End")
    return s


# ============================================================
# Build
# ============================================================

def main():
    here = os.path.dirname(os.path.abspath(__file__))
    repo_root = os.path.abspath(os.path.join(here, "..", ".."))
    video = os.path.join(repo_root, "deliverables", "jury-deck",
                         "assets", "media", "demo_detection.mp4")
    out = os.path.join(here, "HARDWARE_III_jury.pptx")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_00_title(prs)
    slide_01_context_problem(prs)
    slide_02_context_proposition(prs, video)
    slide_03_experience(prs)
    slide_04_interaction(prs)
    slide_05_logic(prs)
    slide_06_architecture(prs)
    slide_07_vision(prs)
    slide_08_gestures(prs)
    slide_09_touchdesigner(prs)
    slide_10_physical(prs)
    slide_11_slider(prs)
    slide_12_data(prs)
    slide_13_metrics(prs)
    slide_14_layout_empty(prs)
    slide_15_layout_roles(prs)
    slide_16_layout_annotated(prs)
    slide_17_fabrication(prs)
    slide_18_bibliography(prs)
    slide_19_end(prs)

    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
